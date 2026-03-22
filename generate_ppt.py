#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零碳园区用能及产业构成PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
TITLE_COLOR = RGBColor(0, 51, 102)  # 深蓝色
ACCENT_COLOR = RGBColor(0, 120, 212)  # 蓝色
SUCCESS_COLOR = RGBColor(34, 139, 34)  # 绿色

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """添加章节页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(4), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_with_chart(prs, title, data_dict, left=1, top=1.5, width=5, height=4):
    """添加带条形图的内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 绘制条形图
    y_pos = top
    bar_height = height / len(data_dict)
    
    for label, value in data_dict.items():
        # 标签
        label_box = slide.shapes.add_textbox(Inches(left), Inches(y_pos), Inches(2), Inches(bar_height*0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.RIGHT
        
        # 条形背景
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 2.2), Inches(y_pos + bar_height*0.2),
            Inches(width), Inches(bar_height * 0.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
        bg_shape.line.fill.background()
        
        # 条形（按比例）
        bar_width = (width * float(value.replace('%', '')) / 100)
        if bar_width > 0.1:
            bar_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 2.2), Inches(y_pos + bar_height*0.2),
                Inches(bar_width), Inches(bar_height * 0.5)
            )
            bar_shape.fill.solid()
            bar_shape.fill.fore_color.rgb = ACCENT_COLOR
            bar_shape.line.fill.background()
        
        # 数值
        value_box = slide.shapes.add_textbox(Inches(left + 2.2 + width + 0.2), Inches(y_pos), Inches(1), Inches(bar_height*0.6))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(14)
        p.font.bold = True
        
        y_pos += bar_height
    
    return slide

def add_two_column_comparison(prs, title, left_data, right_data, left_title, right_title):
    """添加两列对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 左侧标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # 左侧内容
    y_pos = 1.8
    for label, value in left_data.items():
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(5.5), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"  {label}: {value}"
        p.font.size = Pt(14)
        y_pos += 0.45
    
    # 右侧标题
    right_title_box = slide.shapes.add_textbox(Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_COLOR
    
    # 右侧内容
    y_pos = 1.8
    for label, value in right_data.items():
        text_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.5), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"  {label}: {value}"
        p.font.size = Pt(14)
        y_pos += 0.45
    
    return slide

def add_solution_slide(prs, title, solutions):
    """添加解决方案页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 解决方案列表
    y_pos = 1.3
    for i, (category, items) in enumerate(solutions.items()):
        # 类别标题
        cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.5))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"▸ {category}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        y_pos += 0.5
        
        # 详细内容
        for item in items:
            item_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos), Inches(11.5), Inches(0.45))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"  • {item}"
            p.font.size = Pt(14)
            y_pos += 0.4
        
        y_pos += 0.3
    
    return slide

# ========== 生成PPT ==========

# 1. 标题页
add_title_slide(prs, "零碳园区用能及产业构成", "工业节能解决方案")

# 2. 目录
add_section_slide(prs, "目录")

# 3. 汽车制造零碳园区 - 主要工艺用能
add_content_with_chart(prs, "汽车制造零碳园区 - 主要工艺用能占比", {
    "热处理": "39%",
    "涂装": "30%",
    "铸造": "18%",
    "机加工": "13%"
})

# 4. 汽车制造零碳园区 - 主要设备用能
add_content_with_chart(prs, "汽车制造零碳园区 - 主要设备用能占比", {
    "工业炉窑": "49%",
    "空压机": "23%",
    "电机": "15%",
    "空调系统": "13%"
})

# 5. 汽车制造零碳园区 - 公用介质用能
add_content_with_chart(prs, "汽车制造零碳园区 - 公用介质用能占比", {
    "电力": "80%",
    "天然气": "15%",
    "蒸汽": "5%"
})

# 6. 矿山零碳园区 - 主要工艺用能
add_content_with_chart(prs, "矿山零碳园区 - 主要工艺用能占比", {
    "采矿": "41%",
    "选矿": "38%",
    "提升运输": "21%"
})

# 7. 矿山零碳园区 - 主要设备用能
add_content_with_chart(prs, "矿山零碳园区 - 主要设备用能占比", {
    "提升机": "30%",
    "空压机": "25%",
    "水泵": "18%",
    "通风机": "15%",
    "电机": "12%"
})

# 8. 矿山零碳园区 - 公用介质用能
add_content_with_chart(prs, "矿山零碳园区 - 公用介质用能占比", {
    "电力": "91%",
    "柴油": "7%",
    "天然气": "2%"
})

# 9. 行业对比
add_two_column_comparison(prs, "行业用能特征对比", 
    {
        "主导能源": "电力(80%)",
        "最高能耗设备": "工业炉窑(49%)",
        "主要工艺": "热处理(39%)",
        "节能重点": "加热设备"
    },
    {
        "主导能源": "电力(91%)",
        "最高能耗设备": "提升机(30%)",
        "主要工艺": "采矿(41%)",
        "节能重点": "动力设备"
    },
    "汽车制造", "矿山"
)

# 10. 工业节能解决方案
add_solution_slide(prs, "工业节能解决方案", {
    "通用节能措施": [
        "能源管理系统(EMS)实时监测",
        "设备变频改造(VFD)",
        "余热回收利用",
        "智能照明控制系统"
    ],
    "汽车制造专项": [
        "工业炉窑蓄热式燃烧技术",
        "涂装废气余热回收",
        "空压机群联控系统",
        "高效热处理工艺优化"
    ],
    "矿山专项": [
        "提升机能量回馈系统",
        "矿井通风智能控制",
        "水泵高效运行优化",
        "电动设备替代柴油"
    ]
})

# 11. 总结页
add_section_slide(prs, "谢谢")

# 保存文件
output_path = os.path.expanduser("~/.openclaw/workspace/零碳园区用能分析.pptx")
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
