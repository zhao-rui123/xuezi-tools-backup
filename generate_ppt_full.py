#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零碳园区解决方案PPT - 按模板格式
包含：项目概述、用能特征分析、技术方案设计、碳排放特征、零碳路径规划、效益分析
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0, 51, 102)
ACCENT_COLOR = RGBColor(0, 120, 212)
ORANGE_COLOR = RGBColor(255, 140, 0)
GREEN_COLOR = RGBColor(34, 139, 34)
LIGHT_BLUE = RGBColor(200, 220, 240)

def add_cover_slide(prs, title, subtitle, industry):
    """封面页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部色带
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    # 标题背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), prs.slide_width, Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 行业标签
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.6))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📋 {industry}"
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, section_title, color=None):
    """章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg_color = color if color else TITLE_COLOR
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    
    # 章节号
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(3.5), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(5.2), Inches(3), Inches(7.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = bg_color
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_slide(prs, title, content_items, icon="▸"):
    """内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 内容区域 - 两列布局
    y_pos = 1.3
    col1_items = content_items[:len(content_items)//2 + len(content_items)%2]
    col2_items = content_items[len(content_items)//2 + len(content_items)%2:]
    
    x_col1 = 0.5
    x_col2 = 7.0
    
    for i, item in enumerate(col1_items):
        # 项目符号
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_col1), Inches(y_pos + 0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_COLOR
        dot.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(x_col1 + 0.2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        y_pos += 0.5
    
    y_pos = 1.3
    for i, item in enumerate(col2_items):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_col2), Inches(y_pos + 0.08), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_COLOR
        dot.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(x_col2 + 0.2), Inches(y_pos), Inches(6), Inches(0.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        y_pos += 0.5
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, left_color=None, right_color=None):
    """双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    lc = left_color if left_color else ACCENT_COLOR
    rc = right_color if right_color else GREEN_COLOR
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"◆ {left_title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = lc
    
    # 左栏内容
    y_pos = 1.8
    for item in left_items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos + 0.08), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = lc
        dot.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(5.5), Inches(0.45))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        y_pos += 0.45
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"◆ {right_title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = rc
    
    # 右栏内容
    y_pos = 1.8
    for item in right_items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(y_pos + 0.08), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rc
        dot.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.5), Inches(0.45))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        y_pos += 0.45
    
    return slide

def add_table_slide(prs, title, headers, data):
    """表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 表格
    num_rows = len(data) + 1
    table = slide.shapes.add_table(num_rows, len(headers), Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.5)).table
    
    # 设置列宽
    col_widths = [2.5, 3, 3.5, 3.7]
    for i, w in enumerate(col_widths[:len(headers)]):
        table.columns[i].width = Inches(w)
    
    # 表头
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # 数据行
    colors = [RGBColor(240, 248, 255), RGBColor(255, 255, 255)]
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.rows[row_idx+1].cells[col_idx]
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors[row_idx % 2]
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    return slide

def add_timeline_slide(prs, title, phases):
    """时间轴/路径页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 时间轴线
    y_line = Inches(3.5)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_line, Inches(11.7), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # 各阶段
    phase_width = 12.0 / len(phases)
    colors_phase = [ACCENT_COLOR, ORANGE_COLOR, GREEN_COLOR]
    
    for i, (phase_name, goals) in enumerate(phases):
        x_pos = 0.8 + i * phase_width
        
        # 节点圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + phase_width/2 - 0.15), Inches(3.35), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors_phase[i % 3]
        dot.line.fill.background()
        
        # 阶段名称
        name_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.2), Inches(phase_width), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors_phase[i % 3]
        p.alignment = PP_ALIGN.CENTER
        
        # 目标内容
        y_goal = 1.8
        for goal in goals:
            goal_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(y_goal), Inches(phase_width - 0.2), Inches(0.4))
            tf = goal_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {goal}"
            p.font.size = Pt(10)
            y_goal += 0.4
    
    return slide

def add_benefit_slide(prs, title, benefits):
    """效益分析页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 三列效益
    col_width = 4.0
    colors = [ACCENT_COLOR, ORANGE_COLOR, GREEN_COLOR]
    icons = ["💰", "🏭", "🌱"]
    
    for i, (benefit_type, items) in enumerate(benefits.items()):
        x_pos = 0.5 + i * col_width
        
        # 顶部图标和标题
        icon_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.2), Inches(col_width), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icons[i]} {benefit_type}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER
        
        # 分隔线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos + 0.5), Inches(1.9), Inches(col_width - 1), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = colors[i]
        line.line.fill.background()
        
        # 内容
        y_pos = 2.1
        for item in items:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 0.2), Inches(y_pos + 0.08), Inches(0.1), Inches(0.1))
            dot.fill.solid()
            dot.fill.fore_color.rgb = colors[i]
            dot.line.fill.background()
            
            text_box = slide.shapes.add_textbox(Inches(x_pos + 0.4), Inches(y_pos), Inches(col_width - 0.5), Inches(0.45))
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(12)
            y_pos += 0.5
    
    return slide

# ========== 汽车制造零碳园区 ==========

# 1. 封面
add_cover_slide(prs, "零碳园区解决方案", "汽车制造行业", "汽车制造零碳园区")

# 2. 目录
add_section_slide(prs, "目录", "CONTENTS", ACCENT_COLOR)

# 3. 项目概述
add_section_slide(prs, "01", "项目概述", ACCENT_COLOR)
add_content_slide(prs, "项目概述 - 零碳园区建设背景与目标", [
    "【建设背景】",
    "• 汽车制造业是我国国民经济支柱产业，能耗占工业总能耗的15%以上",
    "• 汽车制造工艺复杂，包含铸造、热处理、涂装、机加工等高能耗工序",
    "• 碳排放主要来自能源消耗，其中电力占80%、天然气15%、蒸汽5%",
    "",
    "【建设目标】",
    "• 近期目标(2025)：能效提升20%，碳排放强度下降15%",
    "• 中期目标(2030)：绿电占比达到50%，碳排放强度下降40%",
    "• 远期目标(2035)：实现碳中和，建成零碳智慧园区"
])

# 4. 用能特征分析
add_section_slide(prs, "02", "用能特征分析", ACCENT_COLOR)
add_two_column_slide(prs, "用能特征分析 - 行业工艺特点与主要用能设备",
    "行业工艺特点", [
        "铸造：金属熔炼、高温熔化，能耗占比18%",
        "热处理：精确控温、长时间保温，能耗占比39%",
        "涂装：涂料烘干、废气处理，能耗占比30%",
        "机加工：数控机床、精密加工，能耗占比13%",
        "物流：自动输送、仓储管理"
    ],
    "主要用能设备", [
        "工业炉窑：热处理炉、熔炼炉、烘干炉，占比49%",
        "空压机站：气源供应、喷涂驱动，占比23%",
        "各类电机：主轴、泵、风机，占比15%",
        "空调系统：车间温湿度控制，占比13%",
        "其他：照明、办公、动力"
    ],
    ACCENT_COLOR, GREEN_COLOR
)
add_timeline_slide(prs, "用能结构分析", [
    ("电力主导\n80%", ["工业生产主体能源", "覆盖全部生产工序", "峰谷电价管理"]),
    ("工艺用热\n15%", ["天然气锅炉", "热处理/熔炼", "涂装烘干"]),
    ("辅助蒸汽\n5%", ["清洗脱脂", "工艺加热", "冬季供暖"])
])

# 5. 技术方案设计
add_section_slide(prs, "03", "技术方案设计", ACCENT_COLOR)
add_two_column_slide(prs, "技术方案设计 - 能源替代与设备节能",
    "能源替代方案", [
        "光伏发电：厂房屋顶建设分布式光伏",
        "储能系统：配置储能平抑波动",
        "绿电采购：购买风电/光伏绿证",
        "天然气优化：高效燃烧技术"
    ],
    "设备节能方案", [
        "变频改造：电机、水泵、风机变频",
        "工业炉窑：蓄热燃烧、、余热回收",
        "空压站：群控系统、余热利用",
        "智能照明：LED+感应控制"
    ],
    ACCENT_COLOR, GREEN_COLOR
)
add_two_column_slide(prs, "技术方案设计 - 智能化管理与能源结构优化",
    "智能化管理", [
        "能源管理系统(EMS)全覆盖",
        "设备状态在线监测",
        "AI节能优化算法",
        "碳排放实时核算"
    ],
    "能源结构优化", [
        "电力：91%→65%（含绿电）",
        "光伏：新增装机X MWp",
        "储能：配置X MWh",
        "绿证：年购绿电X MWh"
    ],
    ACCENT_COLOR, GREEN_COLOR
)

# 6. 碳排放特征
add_section_slide(prs, "04", "碳排放特征", ACCENT_COLOR)
add_table_slide(prs, "碳排放结构分析", 
    ["排放源", "占比", "主要来源", "减排措施"],
    [
        ["电力消耗", "65%", "外购电力", "绿电替代+节能"],
        ["天然气燃烧", "25%", "工业炉窑", "高效燃烧+余热"],
        ["蒸汽消耗", "8%", "锅炉制备", "余热利用+制度"],
        ["其他", "2%", "柴油等", "电动化替代"]
    ]
)
add_content_slide(prs, "碳排放强度与减排潜力", [
    "【碳排放强度】",
    "• 当前：约 0.8 tCO2/万元产值",
    "• 目标：2030年降至 0.5 tCO2/万元产值",
    "",
    "【减排潜力分析】",
    "• 节能降耗：可减排 35%（工艺优化、设备升级）",
    "• 绿电替代：可减排 30%（光伏+储能+绿证）",
    "• 智能化管理：可减排 10%（精细化管理）",
    "• 剩余排放：通过碳汇/碳交易中和"
])

# 7. 零碳路径规划
add_section_slide(prs, "05", "零碳路径规划", ACCENT_COLOR)
add_timeline_slide(prs, "零碳路径规划 - 近期/中期/远期目标",
    [
        ("近期\n2023-2025", ["能效提升20%", "LED照明全覆盖", "变频改造", "EMS系统上线", "碳强度-15%"]),
        ("中期\n2026-2030", ["绿电占比50%", "光伏+储能", "余热回收", "智慧能源平台", "碳强度-40%"]),
        ("远期\n2031-2035", ["碳中和", "100%绿电", "深度脱碳", "碳汇抵消", "零碳园区"])
    ]
)

# 8. 效益分析
add_section_slide(prs, "06", "效益分析", ACCENT_COLOR)
add_benefit_slide(prs, "效益分析 - 经济/社会/环境效益",
    {
        "经济效益": [
            "年节能收益约X万元",
            "峰谷电价套利收益",
            "设备维护成本降低",
            "碳交易潜在收益"
        ],
        "社会效益": [
            "带动就业X人",
            "提升企业形象",
            "行业示范引领",
            "助力双碳目标"
        ],
        "环境效益": [
            "年减碳X万吨",
            "粉尘减排X吨",
            "NOx减排X吨",
            "助力生态文明"
        ]
    }
)

# ========== 矿山零碳园区 ==========

# 1. 封面
add_cover_slide(prs, "零碳园区解决方案", "矿山行业", "矿山零碳园区")

# 2. 目录
add_section_slide(prs, "目录", "CONTENTS", ORANGE_COLOR)

# 3. 项目概述
add_section_slide(prs, "01", "项目概述", ORANGE_COLOR)
add_content_slide(prs, "项目概述 - 零碳园区建设背景与目标", [
    "【建设背景】",
    "• 矿业是国民经济基础性产业，能源消耗大，碳排放强度高",
    "• 矿山作业涉及采矿、选矿、输送、排水、通风等环节",
    "• 碳排放主要来自电力消耗（提升、排水、通风占70%以上）",
    "",
    "【建设目标】",
    "• 近期目标(2025)：能效提升25%，碳排放强度下降18%",
    "• 中期目标(2030)：绿电占比达到40%，碳排放强度下降45%",
    "• 远期目标(2035)：实现碳中和，建成绿色智慧矿山"
])

# 4. 用能特征分析
add_section_slide(prs, "02", "用能特征分析", ORANGE_COLOR)
add_two_column_slide(prs, "用能特征分析 - 行业工艺特点与主要用能设备",
    "行业工艺特点", [
        "采矿：凿岩爆破、装载运输，能耗占比41%",
        "选矿：破碎筛分、浮选压滤，能耗占比38%",
        "提升运输：井筒提升、胶带输送，能耗占比21%",
        "排水：井下排水、防尘洒水",
        "通风：矿井通风、瓦斯排放"
    ],
    "主要用能设备", [
        "提升机：主井/副井提升，占比30%",
        "空压机：凿岩供气、喷锚，占比25%",
        "水泵：井下排水，占比18%",
        "通风机：矿井通风，占比15%",
        "电机：传送、破碎等，占比12%"
    ],
    ORANGE_COLOR, GREEN_COLOR
)
add_timeline_slide(prs, "用能结构分析",
    [
        ("电力主导\n91%", ["提升机耗电", "排水泵耗电", "通风机耗电", "照明监控"]),
        ("柴油辅助\n7%", ["运输车辆", "工程设备", "备用电源"]),
        ("天然气\n2%", ["锅炉取暖", "员工设施"])
    ]
)

# 5. 技术方案设计
add_section_slide(prs, "03", "技术方案设计", ORANGE_COLOR)
add_two_column_slide(prs, "技术方案设计 - 能源替代与设备节能",
    "能源替代方案", [
        "光伏发电：办公区/生活区屋顶",
        "储能系统：削峰填谷应急备用",
        "电动设备：电动宽体车替代柴油",
        "绿电采购：购买绿证实现低碳"
    ],
    "设备节能方案", [
        "提升机：变频调速+能量回收",
        "空压站：变频+余热利用+按需供气",
        "通风系统：智能调控+变频",
        "排水系统：变频调速+峰谷运行"
    ],
    ORANGE_COLOR, GREEN_COLOR
)
add_two_column_slide(prs, "技术方案设计 - 智能化管理与能源结构优化",
    "智能化管理", [
        "矿山智慧能源平台",
        "设备远程监控诊断",
        "能耗实时监测分析",
        "碳排放统计核算"
    ],
    "能源结构优化", [
        "电力：91%→55%（含绿电）",
        "光伏：新增装机X MWp",
        "储能：配置X MWh",
        "电动化：柴油减少80%"
    ],
    ORANGE_COLOR, GREEN_COLOR
)

# 6. 碳排放特征
add_section_slide(prs, "04", "碳排放特征", ORANGE_COLOR)
add_table_slide(prs, "碳排放结构分析",
    ["排放源", "占比", "主要来源", "减排措施"],
    [
        ["电力消耗", "75%", "提升/排水/通风", "绿电替代+节能"],
        ["柴油消耗", "20%", "运输/工程设备", "电动化替代"],
        ["天然气燃烧", "3%", "锅炉/取暖", "热泵替代"],
        ["其他", "2%", "炸药/耗材等", "工艺优化"]
    ]
)
add_content_slide(prs, "碳排放强度与减排潜力", [
    "【碳排放强度】",
    "• 当前：约 1.2 tCO2/万吨矿产品",
    "• 目标：2030年降至 0.65 tCO2/万吨矿产品",
    "",
    "【减排潜力分析】",
    "• 设备节能：可减排 30%（变频改造+能量回收）",
    "• 绿电替代：可减排 25%（光伏+储能+绿证）",
    "• 电动化替代：可减排 20%（电动车辆替代柴油）",
    "• 智能管控：可减排 10%（按需供能+精细管理）"
])

# 7. 零碳路径规划
add_section_slide(prs, "05", "零碳路径规划", ORANGE_COLOR)
add_timeline_slide(prs, "零碳路径规划 - 近期/中期/远期目标",
    [
        ("近期\n2023-2025", ["能效提升25%", "提升机变频", "通风智能控制", "EMS平台", "碳强度-18%"]),
        ("中期\n2026-2030", ["绿电占比40%", "电动宽体车", "储能系统", "智能矿山", "碳强度-45%"]),
        ("远期\n2031-2035", ["碳中和", "100%电动化", "100%绿电", "智慧零碳", "碳汇抵消"])
    ]
)

# 8. 效益分析
add_section_slide(prs, "06", "效益分析", ORANGE_COLOR)
add_benefit_slide(prs, "效益分析 - 经济/社会/环境效益",
    {
        "经济效益": [
            "年节能收益约X万元",
            "谷电利用降低成本",
            "设备寿命延长",
            "碳交易收益"
        ],
        "社会效益": [
            "改善井下作业环境",
            "提升安全生产水平",
            "行业示范引领",
            "助力地方经济绿色发展"
        ],
        "环境效益": [
            "年减碳X万吨",
            "粉尘减排X吨",
            "噪音污染降低",
            "生态恢复治理"
        ]
    }
)

# 9. 总结页
add_section_slide(prs, "谢谢", "THANK YOU", TITLE_COLOR)

# 保存文件
output_path = os.path.expanduser("~/.openclaw/workspace/零碳园区解决方案_完整版.pptx")
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
