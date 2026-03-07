#!/usr/bin/env python3
"""
格式转换器
支持 Office 文档和 Markdown、PDF 等格式互转
"""

import argparse
import os
from docx import Document
import pandas as pd
from pptx import Presentation


def docx_to_markdown(input_file, output_file):
    """Word 转 Markdown"""
    doc = Document(input_file)
    md_content = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        
        # 判断标题级别
        if para.style.name.startswith('Heading'):
            level = para.style.name.split()[-1]
            md_content.append(f"{'#' * int(level)} {text}")
        else:
            md_content.append(text)
    
    # 表格
    for table in doc.tables:
        md_content.append('\n')
        for i, row in enumerate(table.rows):
            row_text = '| ' + ' | '.join([cell.text for cell in row.cells]) + ' |'
            md_content.append(row_text)
            if i == 0:
                md_content.append('|' + '|'.join(['---'] * len(row.cells)) + '|')
        md_content.append('')
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write('\n\n'.join(md_content))
    
    print(f"转换完成: {output_file}")


def excel_to_csv(input_file, output_dir):
    """Excel 转 CSV"""
    xls = pd.ExcelFile(input_file)
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(input_file, sheet_name=sheet_name)
        output_file = os.path.join(output_dir, f'{sheet_name}.csv')
        df.to_csv(output_file, index=False, encoding='utf-8-sig')
        print(f"生成: {output_file}")


def excel_to_json(input_file, output_file):
    """Excel 转 JSON"""
    df = pd.read_excel(input_file)
    df.to_json(output_file, orient='records', force_ascii=False, indent=2)
    print(f"转换完成: {output_file}")


def pptx_to_text(input_file, output_file):
    """PPT 提取文本"""
    prs = Presentation(input_file)
    text_content = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_content.append(f"\n=== 幻灯片 {i} ===\n")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text)
    
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write('\n'.join(text_content))
    
    print(f"文本提取完成: {output_file}")


def main():
    parser = argparse.ArgumentParser(description='格式转换器')
    parser.add_argument('action', 
                       choices=['docx2md', 'excel2csv', 'excel2json', 'pptx2txt'],
                       help='转换类型')
    parser.add_argument('--input', '-i', required=True, help='输入文件')
    parser.add_argument('--output', '-o', help='输出文件/目录')
    
    args = parser.parse_args()
    
    if args.action == 'docx2md':
        if not args.output:
            args.output = args.input.replace('.docx', '.md')
        docx_to_markdown(args.input, args.output)
    
    elif args.action == 'excel2csv':
        if not args.output:
            args.output = './csv_output'
        excel_to_csv(args.input, args.output)
    
    elif args.action == 'excel2json':
        if not args.output:
            args.output = args.input.replace('.xlsx', '.json')
        excel_to_json(args.input, args.output)
    
    elif args.action == 'pptx2txt':
        if not args.output:
            args.output = args.input.replace('.pptx', '.txt')
        pptx_to_text(args.input, args.output)


if __name__ == '__main__':
    main()
