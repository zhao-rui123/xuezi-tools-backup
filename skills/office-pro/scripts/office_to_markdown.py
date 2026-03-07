#!/usr/bin/env python3
"""
Office转Markdown工具 - Word/Excel/PPT转Markdown
补充 office-pro 技能包
"""

import os
import sys
from typing import Optional

def docx_to_markdown(docx_path: str, output_path: Optional[str] = None) -> str:
    """Word转Markdown"""
    try:
        from docx import Document
        
        doc = Document(docx_path)
        md_content = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                md_content.append('')
                continue
            
            # 根据样式判断标题级别
            style_name = para.style.name if para.style else 'Normal'
            
            if 'Heading 1' in style_name or style_name == 'Title':
                md_content.append(f"# {text}")
            elif 'Heading 2' in style_name:
                md_content.append(f"## {text}")
            elif 'Heading 3' in style_name:
                md_content.append(f"### {text}")
            elif 'Heading 4' in style_name:
                md_content.append(f"#### {text}")
            else:
                # 普通段落，检查是否列表
                if text.startswith(('•', '-', '*', '·')):
                    md_content.append(f"- {text[1:].strip()}")
                elif len(text) > 0 and text[0].isdigit() and '. ' in text[:4]:
                    md_content.append(text)
                else:
                    md_content.append(text)
        
        # 处理表格
        for table in doc.tables:
            md_content.append('\n')
            # 表头
            header_cells = [cell.text.strip() for cell in table.rows[0].cells]
            md_content.append('| ' + ' | '.join(header_cells) + ' |')
            md_content.append('|' + '|'.join(['---'] * len(header_cells)) + '|')
            
            # 数据行
            for row in table.rows[1:]:
                row_cells = [cell.text.strip() for cell in row.cells]
                md_content.append('| ' + ' | '.join(row_cells) + ' |')
            
            md_content.append('')
        
        result = '\n'.join(md_content)
        
        # 保存
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(result)
            print(f"✅ Word转Markdown完成: {output_path}")
        
        return result
        
    except Exception as e:
        return f"❌ 转换失败: {e}"

def xlsx_to_markdown(xlsx_path: str, output_path: Optional[str] = None, 
                    sheet_index: int = 0) -> str:
    """Excel转Markdown表格"""
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        sheet = wb.worksheets[sheet_index]
        
        md_content = []
        md_content.append(f"## Sheet: {sheet.title}\n")
        
        # 读取所有行
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            return "空表格"
        
        # 转换为字符串并处理None
        def cell_to_str(cell):
            if cell is None:
                return ''
            return str(cell)
        
        # 表头
        header = [cell_to_str(cell) for cell in rows[0]]
        md_content.append('| ' + ' | '.join(header) + ' |')
        md_content.append('|' + '|'.join(['---'] * len(header)) + '|')
        
        # 数据
        for row in rows[1:]:
            row_data = [cell_to_str(cell) for cell in row]
            md_content.append('| ' + ' | '.join(row_data) + ' |')
        
        result = '\n'.join(md_content)
        
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(result)
            print(f"✅ Excel转Markdown完成: {output_path}")
        
        return result
        
    except Exception as e:
        return f"❌ 转换失败: {e}"

def pptx_to_markdown(pptx_path: str, output_path: Optional[str] = None) -> str:
    """PPT转Markdown大纲"""
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        md_content = []
        
        md_content.append(f"# {pptx_path.split('/')[-1]}\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            md_content.append(f"\n## Slide {slide_num}\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # 判断是否为标题（第一个文本框）
                    if shape == slide.shapes[0]:
                        md_content.append(f"### {text}\n")
                    else:
                        md_content.append(f"{text}\n")
        
        result = '\n'.join(md_content)
        
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(result)
            print(f"✅ PPT转Markdown完成: {output_path}")
        
        return result
        
    except Exception as e:
        return f"❌ 转换失败: {e}"

def convert_office_to_md(input_path: str, output_path: Optional[str] = None) -> str:
    """
    自动识别Office文件类型并转Markdown
    
    Args:
        input_path: Office文件路径 (.docx/.xlsx/.pptx)
        output_path: 输出Markdown路径（可选）
    """
    if not os.path.exists(input_path):
        return f"❌ 文件不存在: {input_path}"
    
    # 自动判断类型
    ext = os.path.splitext(input_path)[1].lower()
    
    if output_path is None:
        output_path = input_path.replace(ext, '.md')
    
    if ext == '.docx':
        return docx_to_markdown(input_path, output_path)
    elif ext == '.xlsx':
        return xlsx_to_markdown(input_path, output_path)
    elif ext == '.pptx':
        return pptx_to_markdown(input_path, output_path)
    else:
        return f"❌ 不支持的文件格式: {ext}"

def batch_convert(input_dir: str, output_dir: str):
    """批量转换目录下所有Office文件"""
    import glob
    
    os.makedirs(output_dir, exist_ok=True)
    
    patterns = ['*.docx', '*.xlsx', '*.pptx']
    files = []
    for pattern in patterns:
        files.extend(glob.glob(os.path.join(input_dir, pattern)))
    
    print(f"📁 找到 {len(files)} 个Office文件\n")
    
    for file_path in files:
        filename = os.path.basename(file_path)
        name_without_ext = os.path.splitext(filename)[0]
        output_path = os.path.join(output_dir, f"{name_without_ext}.md")
        
        print(f"🔄 转换: {filename}")
        result = convert_office_to_md(file_path, output_path)
        if result.startswith("❌"):
            print(f"   {result}")
    
    print(f"\n✅ 批量转换完成，输出目录: {output_dir}")

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Office转Markdown工具')
    parser.add_argument('input', help='输入文件或目录')
    parser.add_argument('-o', '--output', help='输出文件或目录')
    parser.add_argument('-b', '--batch', action='store_true', help='批量处理目录')
    
    args = parser.parse_args()
    
    if args.batch:
        output_dir = args.output or './markdown_output'
        batch_convert(args.input, output_dir)
    else:
        result = convert_office_to_md(args.input, args.output)
        if not args.output:
            print(result)
