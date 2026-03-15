#!/usr/bin/env python3
"""
PowerPoint 演示文稿生成器
"""

import argparse
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_presentation(data, output_file, template_path=None):
    """创建 PPT 演示文稿"""
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
    
    # 标题页
    title_slide_layout = prs.slide_layouts[0]  # 标题页布局
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data.get('title', '演示文稿')
    subtitle.text = data.get('subtitle', '')
    
    # 内容页
    for slide_data in data.get('slides', []):
        bullet_slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_data.get('title', '')
        tf = body_shape.text_frame
        
        for item in slide_data.get('content', []):
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
    
    # 图表页（如果有数据）
    for chart_data in data.get('charts', []):
        blank_slide_layout = prs.slide_layouts[5]  # 空白布局
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_box.text_frame.text = chart_data.get('title', '图表')
        
        # 这里可以添加图表（需要更复杂的实现）
        # 简化版：添加表格数据
        rows = len(chart_data.get('data', [])) + 1
        cols = len(chart_data.get('columns', []))
        
        if rows > 1 and cols > 0:
            table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table
            
            # 表头
            for i, col in enumerate(chart_data.get('columns', [])):
                table.cell(0, i).text = col
            
            # 数据
            for i, row_data in enumerate(chart_data.get('data', [])):
                for j, value in enumerate(row_data):
                    table.cell(i + 1, j).text = str(value)
    
    prs.save(output_file)
    print(f"PPT 已生成: {output_file}")


def main():
    parser = argparse.ArgumentParser(description='PPT 生成器')
    parser.add_argument('--data', '-d', required=True, help='数据 JSON 文件')
    parser.add_argument('--output', '-o', required=True, help='输出 PPTX 文件')
    parser.add_argument('--template', '-t', help='模板文件')
    
    args = parser.parse_args()
    
    with open(args.data, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    create_presentation(data, args.output, args.template)


if __name__ == '__main__':
    main()
