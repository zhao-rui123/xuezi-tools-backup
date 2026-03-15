#!/usr/bin/env python3
"""
文档排版优化和精简模块
支持Word、Excel、PPT的格式优化和内容精简
"""

import os
import re
import json
from typing import Dict, List, Optional, Union, Any, Tuple
from dataclasses import dataclass
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.table import WD_TABLE_ALIGNMENT
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN


@dataclass
class OptimizationResult:
    success: bool
    issues_found: List[str]
    issues_fixed: List[str]
    warnings: List[str]
    output_path: str


class DocumentOptimizer:
    """Word文档优化器"""

    DEFAULT_FONT = "宋体"
    DEFAULT_SIZE = 12
    TITLE_SIZE = 22
    HEADING1_SIZE = 16
    HEADING2_SIZE = 14

    def __init__(self, input_path: str):
        self.input_path = input_path
        self.doc = Document(input_path)
        self.issues: List[str] = []
        self.fixed: List[str] = []
        self.warnings: List[str] = []

    def analyze_document(self) -> Dict:
        issues = {
            "format_issues": [],
            "content_issues": [],
            "structure_issues": [],
            "style_issues": []
        }

        for i, para in enumerate(self.doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue

            if len(text) > 500:
                issues["content_issues"].append(f"段落{i+1}过长（{len(text)}字符）")

            if para.style.name.startswith('Normal') and text.startswith('#'):
                issues["format_issues"].append(f"段落{i+1}可能使用了Markdown格式")

            if para.alignment is None:
                issues["format_issues"].append(f"段落{i+1}对齐方式未设置")

        for table in self.doc.tables:
            if not table.style or table.style.name == 'Normal':
                issues["style_issues"].append("表格缺少样式")

        return {
            "total_paragraphs": len([p for p in self.doc.paragraphs if p.text.strip()]),
            "total_tables": len(self.doc.tables),
            "total_images": len(self.doc.inline_shapes),
            "issues": issues
        }

    def fix_formatting(self) -> 'DocumentOptimizer':
        for para in self.doc.paragraphs:
            if not para.text.strip():
                continue

            for run in para.runs:
                if not run.font.name:
                    run.font.name = self.DEFAULT_FONT
                if not run.font.size:
                    run.font.size = Pt(self.DEFAULT_SIZE)

            if para.alignment is None or para.alignment == WD_ALIGN_PARAGRAPH.LEFT:
                para.alignment = WD_ALIGN_PARAGRAPH.LEFT

        self.fixed.append("统一字体和字号")
        return self

    def fix_spacing(self, before: int = 6, after: int = 6, line_space: float = 1.5) -> 'DocumentOptimizer':
        for para in self.doc.paragraphs:
            if para.text.strip():
                para.paragraph_format.space_before = Pt(before)
                para.paragraph_format.space_after = Pt(after)
                para.paragraph_format.line_spacing = line_space

        self.fixed.append(f"设置段落间距（段前{before}pt，段后{after}pt，行距{line_space}）")
        return self

    def normalize_headings(self) -> 'DocumentOptimizer':
        heading_mapping = {
            '标题 1': 'Heading 1',
            '标题 2': 'Heading 2',
            '标题 3': 'Heading 3',
            '标题 4': 'Heading 4',
            '标题 5': 'Heading 5',
        }

        for para in self.doc.paragraphs:
            style_name = para.style.name
            if style_name in heading_mapping:
                para.style = heading_mapping[style_name]

        self.fixed.append("规范化标题样式")
        return self

    def fix_tables(self) -> 'DocumentOptimizer':
        for table in self.doc.tables:
            if table.style and table.style.name == 'Normal':
                table.style = 'Light Grid Accent 1'

            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

        self.fixed.append("优化表格样式")
        return self

    def add_page_numbers(self) -> 'DocumentOptimizer':
        section = self.doc.sections[0]
        footer = section.footer
        paragraph = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()

        run = paragraph.add_run()
        run.add_text("第 ")
        
        fldChar = OxmlElement('w:fldChar')
        fldChar.set(qn('w:fldCharType'), 'begin')
        run._element.append(fldChar)
        
        instrText = OxmlElement('w:instrText')
        instrText.set(qn('w:space'), 'preserve')
        instrText.text = 'PAGE'
        run._element.append(instrText)
        
        fldChar = OxmlElement('w:fldChar')
        fldChar.set(qn('w:fldCharType'), 'end')
        run._element.append(fldChar)
        
        run.add_text(" 页，共 ")
        
        fldChar = OxmlElement('w:fldChar')
        fldChar.set(qn('w:fldCharType'), 'begin')
        run._element.append(fldChar)
        
        instrText = OxmlElement('w:instrText')
        instrText.set(qn('w:space'), 'preserve')
        instrText.text = 'NUMPAGES'
        run._element.append(instrText)
        
        fldChar = OxmlElement('w:fldChar')
        fldChar.set(qn('w:fldCharType'), 'end')
        run._element.append(fldChar)
        
        run.add_text(" 页")
        
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self.fixed.append("添加页码")
        return self

    def compress_whitespace(self) -> 'DocumentOptimizer':
        removed_count = 0
        paragraphs = list(self.doc.paragraphs)
        
        for i in range(len(paragraphs) - 1):
            current = paragraphs[i]
            next_para = paragraphs[i + 1]
            
            if not current.text.strip() and not next_para.text.strip():
                p = current._element
                p.getparent().remove(p)
                removed_count += 1

        self.fixed.append(f"删除{removed_count}个空行")
        return self

    def save(self, output_path: str = None) -> str:
        if output_path is None:
            base, ext = os.path.splitext(self.input_path)
            output_path = f"{base}_optimized{ext}"
        
        self.doc.save(output_path)
        return output_path


class ContentSummarizer:
    """内容精简器 - 提取关键信息"""

    def __init__(self, input_path: str):
        self.input_path = input_path
        self.doc = Document(input_path)

    def extract_headings(self) -> List[Dict]:
        headings = []
        for para in self.doc.paragraphs:
            if para.style.name.startswith('Heading'):
                level = int(para.style.name.split()[-1]) if para.style.name[-1].isdigit() else 1
                headings.append({
                    "level": level,
                    "text": para.text.strip()
                })
        return headings

    def extract_key_points(self) -> List[str]:
        key_points = []
        
        for para in self.doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            indicators = ['关键', '重要', '总结', '结论', '建议', '注意', '核心', '主要', '重点']
            if any(ind in text for ind in indicators):
                key_points.append(text)
                continue

            if para.style.name.startswith('Heading') and int(para.style.name[-1] if para.style.name[-1].isdigit() else 1) <= 2:
                key_points.append(text)

        return key_points[:20]

    def summarize(self, max_length: int = 2000) -> str:
        summary_parts = []
        
        headings = self.extract_headings()
        if headings:
            summary_parts.append("## 文档结构\n")
            for h in headings[:10]:
                summary_parts.append(f"{'#' * h['level']} {h['text']}\n")
        
        key_points = self.extract_key_points()
        if key_points:
            summary_parts.append("\n## 关键要点\n")
            for point in key_points[:10]:
                summary_parts.append(f"- {point[:200]}\n")

        summary = ''.join(summary_parts)
        if len(summary) > max_length:
            summary = summary[:max_length] + "..."

        return summary

    def create_summary_document(self, output_path: str = None) -> str:
        if output_path is None:
            base, ext = os.path.splitext(self.input_path)
            output_path = f"{base}_summary{ext}"

        summary_doc = Document()
        summary_doc.add_heading("文档摘要", 0)
        
        headings = self.extract_headings()
        if headings:
            summary_doc.add_heading("文档结构", level=1)
            for h in headings[:15]:
                level = min(h['level'], 3)
                summary_doc.add_heading(h['text'], level=level)

        key_points = self.extract_key_points()
        if key_points:
            summary_doc.add_heading("关键要点", level=1)
            for point in key_points[:15]:
                summary_doc.add_paragraph(point, style='List Bullet')

        summary_doc.add_paragraph(f"\n生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        
        summary_doc.save(output_path)
        return output_path


class ExcelOptimizer:
    """Excel优化器"""

    def __init__(self, input_path: str):
        self.input_path = input_path
        self.workbook = load_workbook(input_path)

    def analyze(self) -> Dict:
        issues = {
            "formatting": [],
            "data": [],
            "structure": []
        }

        for sheet_name in self.workbook.sheetnames:
            ws = self.workbook[sheet_name]
            
            if ws.max_row > 10000:
                issues["data"].append(f"工作表 '{sheet_name}' 有 {ws.max_row} 行，可能需要分表")

            has_style = False
            for row in ws.iter_rows(max_row=5):
                for cell in row:
                    if cell.fill and cell.fill.start_color and cell.fill.start_color.rgb:
                        has_style = True
                        break
            if not has_style:
                issues["formatting"].append(f"工作表 '{sheet_name}' 缺少样式")

            if ws.freeze_panes is None:
                issues["structure"].append(f"工作表 '{sheet_name}' 未冻结首行")

        return {
            "sheets": len(self.workbook.sheetnames),
            "issues": issues
        }

    def apply_professional_style(self, style: str = "blue") -> 'ExcelOptimizer':
        color_schemes = {
            "blue": {
                "header": "4472C4",
                "header_text": "FFFFFF",
                "alt_row": "D6EAF8"
            },
            "green": {
                "header": "27AE60",
                "header_text": "FFFFFF",
                "alt_row": "D5F5E3"
            },
            "professional": {
                "header": "2C3E50",
                "header_text": "FFFFFF",
                "alt_row": "ECF0F1"
            }
        }

        colors = color_schemes.get(style, color_schemes["blue"])

        for sheet_name in self.workbook.sheetnames:
            ws = self.workbook[sheet_name]
            
            header_fill = PatternFill(start_color=colors["header"], 
                                    end_color=colors["header"], 
                                    fill_type="solid")
            header_font = Font(bold=True, color=colors["header_text"], size=11)
            alt_fill = PatternFill(start_color=colors["alt_row"], 
                                  end_color=colors["alt_row"], 
                                  fill_type="solid")
            
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )

            for col in range(1, ws.max_column + 1):
                cell = ws.cell(1, col)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center', vertical='center')
                cell.border = border

            for row in range(2, ws.max_row + 1):
                for col in range(1, ws.max_column + 1):
                    cell = ws.cell(row, col)
                    cell.border = border
                    if row % 2 == 0:
                        cell.fill = alt_fill

        self.workbook._optimized = True
        return self

    def freeze_header_row(self) -> 'ExcelOptimizer':
        for sheet_name in self.workbook.sheetnames:
            ws = self.workbook[sheet_name]
            ws.freeze_panes = 'A2'
        return self

    def auto_adjust_columns(self) -> 'ExcelOptimizer':
        for sheet_name in self.workbook.sheetnames:
            ws = self.workbook[sheet_name]
            
            for column in ws.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width

        return self

    def remove_empty_rows(self) -> 'ExcelOptimizer':
        removed = 0
        for sheet_name in self.workbook.sheetnames:
            ws = self.workbook[sheet_name]
            
            rows_to_delete = []
            for row_idx in range(ws.max_row, 1, -1):
                is_empty = all(ws.cell(row_idx, col).value is None 
                             for col in range(1, ws.max_column + 1))
                if is_empty:
                    rows_to_delete.append(row_idx)
            
            for row_idx in rows_to_delete:
                ws.delete_rows(row_idx)
                removed += 1

        return self

    def save(self, output_path: str = None) -> str:
        if output_path is None:
            base, ext = os.path.splitext(self.input_path)
            output_path = f"{base}_optimized{ext}"
        
        self.workbook.save(output_path)
        return output_path


class PPTOptimizer:
    """PowerPoint优化器"""

    def __init__(self, input_path: str):
        self.input_path = input_path
        self.presentation = Presentation(input_path)

    def analyze(self) -> Dict:
        issues = {
            "formatting": [],
            "design": [],
            "content": []
        }

        for i, slide in enumerate(self.presentation.slides):
            shapes_count = len(shapes := list(slide.shapes))
            
            if shapes_count == 0:
                issues["content"].append(f"幻灯片 {i+1} 为空")
            elif shapes_count > 10:
                issues["content"].append(f"幻灯片 {i+1} 内容过多（{shapes_count}个元素）")

            title_found = any(shape == slide.shapes.title for shape in slide.shapes if hasattr(slide.shapes, 'title'))
            if not title_found and i > 0:
                issues["formatting"].append(f"幻灯片 {i+1} 缺少标题")

        return {
            "total_slides": len(self.presentation.slides),
            "issues": issues
        }

    def apply_theme(self, theme: str = "corporate") -> 'PPTOptimizer':
        themes = {
            "corporate": {
                "primary": PptRGBColor(31, 78, 121),
                "secondary": PptRGBColor(46, 117, 182),
                "accent": PptRGBColor(112, 173, 71)
            },
            "modern": {
                "primary": PptRGBColor(44, 62, 80),
                "secondary": PptRGBColor(52, 73, 94),
                "accent": PptRGBColor(231, 76, 60)
            },
            "elegant": {
                "primary": PptRGBColor(26, 26, 46),
                "secondary": PptRGBColor(22, 33, 62),
                "accent": PptRGBColor(233, 69, 96)
            }
        }

        colors = themes.get(theme, themes["corporate"])

        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if not run.font.color.rgb:
                                run.font.color.rgb = colors["primary"]

        return self

    def normalize_text_sizes(self) -> 'PPTOptimizer':
        for slide in self.presentation.slides:
            title_shape = slide.shapes.title
            if title_shape:
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = PptPt(32)

            for shape in slide.shapes:
                if shape != title_shape and hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size and run.font.size.pt > 18:
                                run.font.size = PptPt(14)

        return self

    def remove_empty_slides(self) -> 'PPTOptimizer':
        slides_to_remove = []
        
        for i, slide in enumerate(self.presentation.slides):
            has_content = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    has_content = True
                    break
            
            if not has_content:
                slides_to_remove.append(slide)

        for slide in reversed(self.presentation.slides._sldIdLst):
            if slide in slides_to_remove:
                self.presentation.slides._sldIdLst.remove(slide)

        return self

    def add_slide_numbers(self) -> 'PPTOptimizer':
        for i, slide in enumerate(self.presentation.slides):
            textbox = slide.shapes.add_textbox(
                PptInches(8.5), PptInches(7), 
                PptInches(1), PptInches(0.3)
            )
            textbox.text_frame.text = f"{i+1}/{len(self.presentation.slides)}"
            textbox.text_frame.paragraphs[0].font.size = PptPt(10)
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        return self

    def save(self, output_path: str = None) -> str:
        if output_path is None:
            base, ext = os.path.splitext(self.input_path)
            output_path = f"{base}_optimized{ext}"
        
        self.presentation.save(output_path)
        return output_path


def optimize_word(input_path: str, output_path: str = None,
                 fix_format: bool = True, fix_spacing: bool = True,
                 fix_tables: bool = True, add_page_numbers: bool = True,
                 compress: bool = True) -> OptimizationResult:
    optimizer = DocumentOptimizer(input_path)
    issues = optimizer.analyze_document()

    if fix_format:
        optimizer.fix_formatting()
    if fix_spacing:
        optimizer.fix_spacing()
    if fix_tables:
        optimizer.fix_tables()
    if add_page_numbers:
        optimizer.add_page_numbers()
    if compress:
        optimizer.compress_whitespace()

    output = optimizer.save(output_path)

    return OptimizationResult(
        success=True,
        issues_found=[f"{k}: {len(v)}" for k, v in issues["issues"].items()],
        issues_fixed=optimizer.fixed,
        warnings=optimizer.warnings,
        output_path=output
    )


def summarize_document(input_path: str, output_path: str = None,
                     extract_key_points: bool = True) -> str:
    if input_path.endswith('.docx'):
        summarizer = ContentSummarizer(input_path)
        return summarizer.create_summary_document(output_path)
    else:
        raise ValueError("仅支持Word文档")


def optimize_excel(input_path: str, output_path: str = None,
                   apply_style: str = "blue", freeze_header: bool = True,
                   auto_adjust: bool = True, remove_empty: bool = True) -> OptimizationResult:
    optimizer = ExcelOptimizer(input_path)
    issues = optimizer.analyze()

    if apply_style:
        optimizer.apply_professional_style(apply_style)
    if freeze_header:
        optimizer.freeze_header_row()
    if auto_adjust:
        optimizer.auto_adjust_columns()
    if remove_empty:
        optimizer.remove_empty_rows()

    output = optimizer.save(output_path)

    return OptimizationResult(
        success=True,
        issues_found=[f"{k}: {len(v)}" for k, v in issues["issues"].items()],
        issues_fixed=["应用样式", "冻结首行", "自动调整列宽", "删除空行"],
        warnings=[],
        output_path=output
    )


def optimize_ppt(input_path: str, output_path: str = None,
                apply_theme: str = "corporate", normalize_text: bool = True,
                remove_empty: bool = True, add_numbers: bool = True) -> OptimizationResult:
    optimizer = PPTOptimizer(input_path)
    issues = optimizer.analyze()

    if apply_theme:
        optimizer.apply_theme(apply_theme)
    if normalize_text:
        optimizer.normalize_text_sizes()
    if remove_empty:
        optimizer.remove_empty_slides()
    if add_numbers:
        optimizer.add_slide_numbers()

    output = optimizer.save(output_path)

    return OptimizationResult(
        success=True,
        issues_found=[f"{k}: {len(v)}" for k, v in issues["issues"].items()],
        issues_fixed=["应用主题", "统一字号", "删除空幻灯片", "添加页码"],
        warnings=[],
        output_path=output
    )


if __name__ == "__main__":
    print("文档优化工具测试")
    
    print("\n1. Word文档分析示例:")
    print("   from core import optimize_word")
    print("   result = optimize_word('input.docx', 'output.docx')")
    print("   print(result.issues_fixed)")
    
    print("\n2. 文档摘要提取:")
    print("   from core import summarize_document")
    print("   summarize_document('input.docx', 'summary.docx')")
    
    print("\n3. Excel优化:")
    print("   from core import optimize_excel")
    print("   optimize_excel('data.xlsx', 'optimized.xlsx', apply_style='blue')")
    
    print("\n4. PPT优化:")
    print("   from core import optimize_ppt")
    print("   optimize_ppt('presentation.pptx', 'optimized.pptx', apply_theme='modern')")
