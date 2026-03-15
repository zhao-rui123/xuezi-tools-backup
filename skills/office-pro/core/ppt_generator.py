#!/usr/bin/env python3
"""
PowerPoint 模板生成器
支持多种专业样式的PPT模板
"""

import os
import json
from typing import Dict, List, Optional, Union, Any
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsmap
import pandas as pd


class PPTStyle:
    STYLES = {
        "corporate": {
            "name": "企业蓝",
            "primary": "1F4E79",
            "secondary": "2E75B6",
            "accent": "70AD47",
            "background": "FFFFFF",
            "text": "1F4E79",
            "title_size": 40,
            "body_size": 18
        },
        "modern": {
            "name": "现代简约",
            "primary": "2C3E50",
            "secondary": "34495E",
            "accent": "E74C3C",
            "background": "FFFFFF",
            "text": "2C3E50",
            "title_size": 44,
            "body_size": 20
        },
        "elegant": {
            "name": "优雅深色",
            "primary": "1A1A2E",
            "secondary": "16213E",
            "accent": "E94560",
            "background": "1A1A2E",
            "text": "FFFFFF",
            "title_size": 42,
            "body_size": 18
        },
        "fresh": {
            "name": "清新绿",
            "primary": "2D5A27",
            "secondary": "4CAF50",
            "accent": "8BC34A",
            "background": "F1F8E9",
            "text": "2D5A27",
            "title_size": 38,
            "body_size": 18
        },
        "warm": {
            "name": "暖色调",
            "primary": "8D4004",
            "secondary": "E67E22",
            "accent": "F39C12",
            "background": "FFF8E1",
            "text": "8D4004",
            "title_size": 40,
            "body_size": 18
        },
        "purple": {
            "name": "商务紫",
            "primary": "4A148C",
            "secondary": "7B1FA2",
            "accent": "AB47BC",
            "background": "F3E5F5",
            "text": "4A148C",
            "title_size": 40,
            "body_size": 18
        },
        "tech": {
            "name": "科技蓝",
            "primary": "0D47A1",
            "secondary": "1976D2",
            "accent": "00BCD4",
            "background": "E3F2FD",
            "text": "0D47A1",
            "title_size": 42,
            "body_size": 18
        }
    }

    @classmethod
    def get_style(cls, style_name: str = "corporate") -> Dict:
        return cls.STYLES.get(style_name, cls.STYLES["corporate"])

    @staticmethod
    def hex_to_rgb(hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


class PPTTemplateGenerator:
    def __init__(self):
        self.presentation = None
        self.style = None
        self.current_slide = None

    def create_presentation(self, style: str = "corporate", width: float = 10, height: float = 7.5):
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(width)
        self.presentation.slide_height = Inches(height)
        self.style = PPTStyle.get_style(style)
        return self

    def add_title_slide(self, title: str, subtitle: str = None):
        layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        title_text_frame = title_shape.text_frame
        for paragraph in title_text_frame.paragraphs:
            paragraph.font.size = Pt(self.style["title_size"])
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPTStyle.hex_to_rgb(self.style["primary"])
            paragraph.alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle

            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(self.style["body_size"])
                paragraph.font.color.rgb = PPTStyle.hex_to_rgb(self.style["secondary"])
                paragraph.alignment = PP_ALIGN.CENTER

        self.current_slide = slide
        return self

    def add_content_slide(self, title: str, bullets: List[str] = None,
                          content: str = None, image_path: str = None):
        layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(self.style["title_size"] - 4)
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPTStyle.hex_to_rgb(self.style["primary"])

        if bullets:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(self.style["body_size"])
                p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["text"])
                p.level = 0

        elif content:
            body_shape = slide.placeholders[1]
            body_shape.text_frame.text = content

            for paragraph in body_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(self.style["body_size"])
                paragraph.font.color.rgb = PPTStyle.hex_to_rgb(self.style["text"])

        elif image_path:
            if os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(1), Inches(2),
                                        width=Inches(8), height=Inches(4.5))

        self.current_slide = slide
        return self

    def add_two_column_slide(self, title: str, left_title: str, left_content: List[str],
                            right_title: str = None, right_content: List[str] = None):
        layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.2), Inches(4.5))
        left_tf = left_box.text_frame
        left_p = left_tf.add_paragraph()
        left_p.text = left_title
        left_p.font.size = Pt(self.style["title_size"] - 8)
        left_p.font.bold = True
        left_p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["secondary"])

        for item in left_content:
            p = left_tf.add_paragraph()
            p.text = item
            p.font.size = Pt(self.style["body_size"] - 2)
            p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["text"])
            p.level = 1

        if right_title and right_content:
            right_box = slide.shapes.add_textbox(Inches(5.3), Inches(2), Inches(4.2), Inches(4.5))
            right_tf = right_box.text_frame
            right_p = right_tf.add_paragraph()
            right_p.text = right_title
            right_p.font.size = Pt(self.style["title_size"] - 8)
            right_p.font.bold = True
            right_p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["secondary"])

            for item in right_content:
                p = right_tf.add_paragraph()
                p.text = item
                p.font.size = Pt(self.style["body_size"] - 2)
                p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["text"])
                p.level = 1

        self.current_slide = slide
        return self

    def add_chart_slide(self, title: str, data: Dict, chart_type: str = "bar"):
        layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(self.style["title_size"] - 4)
        title_p.font.bold = True
        title_p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["primary"])

        rows = len(data.get("data", [])) + 1
        cols = len(data.get("columns", []))

        if rows > 1 and cols > 0:
            table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2),
                                           Inches(8), Inches(4)).table

            for i, col in enumerate(data.get("columns", [])):
                cell = table.cell(0, i)
                cell.text = str(col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PPTStyle.hex_to_rgb(self.style["primary"])
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].font.bold = True

            for i, row_data in enumerate(data.get("data", [])):
                for j, value in enumerate(row_data):
                    cell = table.cell(i + 1, j)
                    cell.text = str(value)

        self.current_slide = slide
        return self

    def add_image_slide(self, title: str, image_path: str, caption: str = None):
        layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(self.style["title_size"] - 4)
        title_p.font.bold = True
        title_p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["primary"])

        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5),
                                    width=Inches(9), height=Inches(5))

            if caption:
                caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
                caption_tf = caption_box.text_frame
                caption_p = caption_tf.paragraphs[0]
                caption_p.text = caption
                caption_p.font.size = Pt(12)
                caption_p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["secondary"])
                caption_p.alignment = PP_ALIGN.CENTER

        self.current_slide = slide
        return self

    def add_closing_slide(self, title: str, contact_info: Dict = None):
        layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(self.style["title_size"])
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPTStyle.hex_to_rgb(self.style["primary"])
            paragraph.alignment = PP_ALIGN.CENTER

        if contact_info:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame

            for key, value in contact_info.items():
                p = tf.add_paragraph()
                p.text = f"{key}: {value}"
                p.font.size = Pt(self.style["body_size"])
                p.font.color.rgb = PPTStyle.hex_to_rgb(self.style["secondary"])
                p.alignment = PP_ALIGN.CENTER

        self.current_slide = slide
        return self

    def add_section_divider(self, title: str):
        layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(self.style["title_size"])
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPTStyle.hex_to_rgb(self.style["accent"])
            paragraph.alignment = PP_ALIGN.CENTER

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.5),
                                       Inches(4), Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PPTStyle.hex_to_rgb(self.style["accent"])

        self.current_slide = slide
        return self

    def add_progress_bar(self, percentage: int):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

        title = slide.shapes.title
        title.text = "进度"

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3),
                                       Inches(8), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PPTStyle.hex_to_rgb(self.style["secondary"])

        progress = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3),
                                         Inches(8 * percentage / 100), Inches(1))
        progress.fill.solid()
        progress.fill.fore_color.rgb = PPTStyle.hex_to_rgb(self.style["accent"])

        text_box = slide.shapes.add_textbox(Inches(4), Inches(3.2), Inches(2), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{percentage}%"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        return self

    def save(self, output_path: str):
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        self.presentation.save(output_path)
        print(f"PPT 已生成: {output_path}")
        return output_path


class PPTTemplateLibrary:
    TEMPLATES = {
        "business_report": {
            "name": "商务报告",
            "style": "corporate",
            "description": "专业的商务报告演示模板",
            "slides": ["封面", "目录", "执行摘要", "数据概览", "分析", "结论", "致谢"]
        },
        "product_launch": {
            "name": "产品发布",
            "style": "modern",
            "description": "产品发布会演示模板",
            "slides": ["封面", "产品介绍", "核心功能", "市场机会", "竞争优势", "定价策略", "总结"]
        },
        "quarterly_review": {
            "name": "季度回顾",
            "style": "tech",
            "description": "季度业务回顾模板",
            "slides": ["封面", "业绩概览", "销售数据", "运营数据", "财务分析", "下季度计划", "Q&A"]
        },
        "project_proposal": {
            "name": "项目提案",
            "style": "elegant",
            "description": "项目提案演示模板",
            "slides": ["封面", "项目背景", "目标与范围", "实施方案", "时间线", "预算", "预期成果", "下一步"]
        },
        "training": {
            "name": "培训课件",
            "style": "fresh",
            "description": "培训演示模板",
            "slides": ["封面", "课程目标", "内容概览", "第一部分", "第二部分", "第三部分", "总结", "问答"]
        },
        "investor_pitch": {
            "name": "投资者演示",
            "style": "warm",
            "description": "投资人路演模板",
            "slides": ["封面", "团队介绍", "问题与解决方案", "市场分析", "商业模式", "竞争优势", "财务预测", "融资需求"]
        }
    }

    @classmethod
    def list_templates(cls) -> List[Dict]:
        return [
            {"id": k, "name": v["name"], "description": v["description"], "style": v["style"]}
            for k, v in cls.TEMPLATES.items()
        ]

    @classmethod
    def create_template(cls, template_id: str, data: Dict = None, output_path: str = None) -> str:
        if template_id not in cls.TEMPLATES:
            raise ValueError(f"Template not found: {template_id}")

        template = cls.TEMPLATES[template_id]
        generator = PPTTemplateGenerator()
        generator.create_presentation(style=template["style"])

        if template_id == "business_report":
            generator.add_title_slide("商务报告", "专业数据驱动决策")
            generator.add_content_slide("目录", ["执行摘要", "数据分析", "结论与建议"])
            generator.add_section_divider("执行摘要")
            generator.add_content_slide("执行摘要", bullets=["本报告分析了...", "关键发现：", "- 发现1", "- 发现2"])
            generator.add_section_divider("数据分析")
            generator.add_content_slide("数据分析", bullets=["数据来源", "分析方法", "主要结果"])
            generator.add_section_divider("结论与建议")
            generator.add_content_slide("结论与建议", bullets=["主要结论", "行动建议", "下一步计划"])
            generator.add_closing_slide("感谢聆听", {"联系人": "example@company.com"})

        elif template_id == "quarterly_review":
            generator.add_title_slide("季度业务回顾", f"Q1 {datetime.now().year}")
            generator.add_content_slide("业绩概览", bullets=["销售业绩", "利润分析", "增长趋势"])
            generator.add_two_column_slide("销售数据", "本月", ["销售额：¥1,000,000", "同比增长：15%"],
                                          "本季度", ["销售额：¥3,000,000", "同比增长：20%"])
            generator.add_content_slide("运营数据", bullets=["客户数量", "服务质量", "运营效率"])
            generator.add_content_slide("财务分析", bullets=["收入结构", "成本分析", "利润空间"])
            generator.add_content_slide("下季度计划", bullets=["目标", "策略", "资源需求"])
            generator.add_closing_slide("Q&A", {"联系": "谢谢"})

        elif template_id == "product_launch":
            generator.add_title_slide("产品发布会", "全新产品震撼登场")
            generator.add_content_slide("产品介绍", bullets=["产品名称", "产品定位", "核心价值"])
            generator.add_content_slide("核心功能", bullets=["功能1", "功能2", "功能3"])
            generator.add_content_slide("市场机会", bullets=["市场规模", "目标用户", "增长潜力"])
            generator.add_content_slide("竞争优势", bullets=["独特卖点", "技术优势", "品牌优势"])
            generator.add_content_slide("定价策略", bullets=["定价方案", "销售渠道", "推广计划"])
            generator.add_closing_slide("谢谢", {"官网": "www.example.com"})

        else:
            generator.add_title_slide(template["name"], template["description"])
            for slide_name in template["slides"][1:-1]:
                generator.add_content_slide(slide_name, bullets=[f"{slide_name}内容占位符"])
            generator.add_closing_slide("感谢")

        if output_path is None:
            output_path = f"{template_id}_{datetime.now().strftime('%Y%m%d')}.pptx"

        generator.save(output_path)
        return output_path


def create_ppt_with_template(template_id: str, data: Dict = None,
                             output_path: str = None) -> str:
    return PPTTemplateLibrary.create_template(template_id, data, output_path)


def create_professional_ppt(output_path: str, title: str,
                           style: str = "corporate",
                           slides: List[Dict] = None) -> str:
    generator = PPTTemplateGenerator()
    generator.create_presentation(style=style)

    generator.add_title_slide(title, "专业演示")

    if slides:
        for slide in slides:
            slide_type = slide.get("type", "content")

            if slide_type == "title":
                generator.add_title_slide(slide.get("title", ""), slide.get("subtitle"))
            elif slide_type == "content":
                generator.add_content_slide(slide.get("title", ""), slide.get("bullets"))
            elif slide_type == "two_column":
                generator.add_two_column_slide(
                    slide.get("title", ""),
                    slide.get("left_title", ""), slide.get("left_content", []),
                    slide.get("right_title"), slide.get("right_content")
                )
            elif slide_type == "chart":
                generator.add_chart_slide(slide.get("title", ""), slide.get("data", {}))
            elif slide_type == "image":
                generator.add_image_slide(slide.get("title", ""), slide.get("image_path", ""))
            elif slide_type == "section":
                generator.add_section_divider(slide.get("title", ""))
            elif slide_type == "closing":
                generator.add_closing_slide(slide.get("title", "感谢"), slide.get("contact"))

    generator.save(output_path)
    return output_path


if __name__ == "__main__":
    print("创建专业PPT...")

    create_professional_ppt(
        "professional_presentation.pptx",
        title="年度总结报告",
        style="corporate",
        slides=[
            {"type": "title", "title": "年度总结报告", "subtitle": "2024年业务回顾与展望"},
            {"type": "content", "title": "目录", "bullets": ["业绩概览", "数据分析", "未来展望"]},
            {"type": "section", "title": "业绩概览"},
            {"type": "content", "title": "主要成就", "bullets": ["销售增长20%", "客户满意度提升", "新产品发布"]},
            {"type": "two_column", "title": "数据分析",
             "left_title": "销售", "left_content": ["¥1000万", "增长15%"],
             "right_title": "利润", "right_content": ["¥200万", "增长25%"]},
            {"type": "section", "title": "未来展望"},
            {"type": "content", "title": "2025年计划", "bullets": ["扩展市场", "产品创新", "团队建设"]},
            {"type": "closing", "title": "感谢聆听", "contact": {"联系": "info@company.com"}}
        ]
    )

    print("\n使用PPT模板库...")
    for tmpl in PPTTemplateLibrary.list_templates():
        print(f"  - {tmpl['name']}: {tmpl['description']}")

    PPTTemplateLibrary.create_template("quarterly_review", None, "quarterly_template.pptx")
    print("\nPPT模板生成完成!")
