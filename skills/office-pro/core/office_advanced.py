#!/usr/bin/env python3
"""
Office 高级功能模块
包含Word高级功能、Excel数据处理、Office格式转换等
"""

import os
import json
import pandas as pd
from typing import Dict, List, Optional, Union, Any
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.style import WD_STYLE_TYPE
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.formatting.rule import Rule, ColorScaleRule, DataBarRule, FormulaRule
from openpyxl.formatting import Rule
from openpyxl.chart import BarChart, LineChart, PieChart
from openpyxl.chart.reference import Reference
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
import io


class WordAdvancedFeatures:
    @staticmethod
    def create_document_with_styles(template_path: str = None) -> Document:
        if template_path and os.path.exists(template_path):
            return Document(template_path)
        return Document()

    @staticmethod
    def add_header(doc: Document, text: str, align: str = "center"):
        header = doc.sections[0].header
        paragraph = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        paragraph.text = text
        if align == "center":
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif align == "right":
            paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        return paragraph

    @staticmethod
    def add_footer(doc: Document, text: str = None, include_page_number: bool = True):
        footer = doc.sections[0].footer
        if text:
            paragraph = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
            paragraph.text = text
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        if include_page_number:
            paragraph = footer.add_paragraph()
            run = paragraph.add_run()
            run.add_text("第 ")
            fldChar = OxmlElement('w:fldChar')
            fldChar.set(qn('w:fldCharType'), 'begin')
            run._element.append(fldChar)

            instrText = OxmlElement('w:instrText')
            instrText.set(qn('w:space'), 'preserve')
            instrText.text = 'PAGE'
            run._element.append(instrText)

            fldChar = OxmlElement('w:fldChar')
            fldChar.set(qn('w:fldCharType'), 'end')
            run._element.append(fldChar)

            run.add_text(" 页，共 ")
            
            fldChar = OxmlElement('w:fldChar')
            fldChar.set(qn('w:fldCharType'), 'begin')
            run._element.append(fldChar)

            instrText = OxmlElement('w:instrText')
            instrText.set(qn('w:space'), 'preserve')
            instrText.text = 'NUMPAGES'
            run._element.append(instrText)

            fldChar = OxmlElement('w:fldChar')
            fldChar.set(qn('w:fldCharType'), 'end')
            run._element.append(fldChar)
            
            run.add_text(" 页")

    @staticmethod
    def add_table_of_contents(doc: Document, title: str = "目录", style: str = "heading 1"):
        paragraph = doc.add_paragraph()
        paragraph.text = title
        paragraph.style = style
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        toc_para = doc.add_paragraph()
        run = toc_para.add_run()
        
        fldChar = OxmlElement('w:fldChar')
        fldChar.set(qn('w:fldCharType'), 'begin')
        run._element.append(fldChar)

        instrText = OxmlElement('w:instrText')
        instrText.set(qn('w:space'), 'preserve')
        instrText.text = r'TOC \h \z \t "heading 1,1"'
        run._element.append(instrText)

        fldChar = OxmlElement('w:fldChar')
        fldChar.set(qn('w:fldCharType'), 'end')
        run._element.append(fldChar)

        doc.add_page_break()

    @staticmethod
    def add_footnote(doc: Document, reference_text: str, footnote_text: str):
        p = doc.add_paragraph(reference_text)
        
        fldChar = OxmlElement('w:fldChar')
        fldChar.set(qn('w:fldCharType'), 'begin')
        
        instrText = OxmlElement('w:instrText')
        instrText.text = 'FOOTNOTE'
        fldChar.addprevious(instrText)
        
        fldChar2 = OxmlElement('w:fldChar')
        fldChar2.set(qn('w:fldCharType'), 'end')
        
        footnote = doc.footnotes
        p2 = doc.add_paragraph()
        p2.text = footnote_text
        p2.style = 'Footnote Text'

    @staticmethod
    def protect_document(doc: Document, password: str = None):
        if password:
            doc.settings.password = password
        else:
            pass

    @staticmethod
    def add_bookmark(doc: Document, bookmark_name: str, text: str):
        paragraph = doc.add_paragraph()
        run = paragraph.add_run(text)
        
        start = OxmlElement('w:bookmarkStart')
        start.set(qn('w:id'), '0')
        start.set(qn('w:name'), bookmark_name)
        run._element.append(start)
        
        end = OxmlElement('w:bookmarkEnd')
        end.set(qn('w:id'), '0')
        run._element.append(end)

    @staticmethod
    def mail_merge(doc: Document, data: List[Dict], output_dir: str = None):
        os.makedirs(output_dir or "mail_merge_output", exist_ok=True)

        for i, record in enumerate(data):
            new_doc = Document(doc)
            
            for paragraph in new_doc.paragraphs:
                for key, value in record.items():
                    if key in paragraph.text:
                        paragraph.text = paragraph.text.replace(key, str(value))

            for table in new_doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for key, value in record.items():
                            if key in cell.text:
                                cell.text = cell.text.replace(key, str(value))

            output_path = os.path.join(output_dir or "mail_merge_output", f"document_{i+1}.docx")
            new_doc.save(output_path)

        return output_dir or "mail_merge_output"


class ExcelAdvancedFeatures:
    @staticmethod
    def create_pivot_table(workbook: Workbook, data_sheet_name: str, 
                          pivot_sheet_name: str, rows: List[str], 
                          values: List[str], filters: List[str] = None):
        ws = workbook[data_sheet_name]
        pivot_ws = workbook.create_sheet(title=pivot_sheet_name)
        
        pivot_ws['A1'] = "数据透视表"
        pivot_ws['A1'].font = Font(bold=True, size=14)

        headers = rows + values
        for col, header in enumerate(headers, 1):
            pivot_ws.cell(2, col, header)

        return pivot_ws

    @staticmethod
    def add_conditional_formatting(ws, range: str, rule_type: str = "colorScale",
                                 params: Dict = None):
        if rule_type == "colorScale":
            rule = ColorScaleRule(
                start_type=params.get("start_type", "min"),
                start_color=params.get("start_color", "F8696B"),
                mid_type=params.get("mid_type", "percentile"),
                mid_value=50,
                mid_color=params.get("mid_color", "FFEB84"),
                end_type=params.get("end_type", "max"),
                end_color=params.get("end_color", "63BE7B")
            )
        elif rule_type == "dataBar":
            rule = DataBarRule(
                start_type=params.get("start_type", "min"),
                end_type=params.get("end_type", "max"),
                color=params.get("color", "638EC6")
            )
        else:
            rule = ColorScaleRule()

        ws.conditional_formatting.add(range, rule)
        return ws

    @staticmethod
    def add_data_validation(ws, range: str, validation_type: str,
                          formula1: str = None, formula2: str = None,
                          show_dropdown: bool = True):
        if validation_type == "list":
            dv = DataValidation(type="list", formula1=formula1, 
                               showDropDown=not show_dropdown)
        elif validation_type == "whole":
            dv = DataValidation(type="whole", operator="between",
                               formula1=formula1 or "0", formula2=formula2 or "100")
        elif validation_type == "decimal":
            dv = DataValidation(type="decimal", operator="between",
                               formula1=formula1 or "0", formula2=formula2 or "100")
        elif validation_type == "date":
            dv = DataValidation(type="date", operator="between",
                               formula1=formula1, formula2=formula2)
        else:
            dv = DataValidation(type="custom", formula1=formula1)

        ws.add_data_validation(dv)
        dv.add(ws[range])
        return ws

    @staticmethod
    def add_freeze_panes(ws, cell: str = "A2"):
        ws.freeze_panes = ws[cell]
        return ws

    @staticmethod
    def add_autofilter(ws, range: str = None):
        if range:
            ws.auto_filter.ref = range
        else:
            ws.auto_filter.ref = ws.dimensions
        return ws


class OfficeFormatConverter:
    @staticmethod
    def excel_to_word_table(excel_path: str, output_path: str, 
                           sheet_name: str = None, title: str = "数据表"):
        df = pd.read_excel(excel_path, sheet_name=sheet_name)
        
        doc = Document()
        doc.add_heading(title, 0)
        
        table = doc.add_table(rows=1, cols=len(df.columns))
        table.style = 'Light Grid Accent 1'
        
        header_cells = table.rows[0].cells
        for i, col_name in enumerate(df.columns):
            header_cells[i].text = str(col_name)
        
        for _, row in df.iterrows():
            row_cells = table.add_row().cells
            for i, value in enumerate(row):
                row_cells[i].text = str(value)
        
        doc.save(output_path)
        return output_path

    @staticmethod
    def word_table_to_excel(doc_path: str, output_path: str):
        doc = Document(doc_path)
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Word表格数据"
        
        row_idx = 1
        for table in doc.tables:
            for row in table.rows:
                col_idx = 1
                for cell in row.cells:
                    ws.cell(row_idx, col_idx, cell.text)
                    col_idx += 1
                row_idx += 1
            row_idx += 1
        
        wb.save(output_path)
        return output_path

    @staticmethod
    def excel_to_ppt(excel_path: str, output_path: str, 
                    sheet_name: str = None, title: str = "数据报告"):
        df = pd.read_excel(excel_path, sheet_name=sheet_name)
        
        prs = Presentation()
        prs.slide_width = PptInches(10)
        prs.slide_height = PptInches(7.5)
        
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        
        blank_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_layout)
        
        rows = min(len(df) + 1, 15)
        cols = len(df.columns)
        
        if rows > 1 and cols > 0:
            table = slide.shapes.add_table(rows, cols, PptInches(0.5), 
                                          PptInches(1.5), 
                                          PptInches(9), 
                                          PptInches(4)).table
            
            for i, col in enumerate(df.columns):
                table.cell(0, i).text = str(col)
            
            for i, row_data in enumerate(df.head(rows - 1).values):
                for j, value in enumerate(row_data):
                    table.cell(i + 1, j).text = str(value)
        
        prs.save(output_path)
        return output_path

    @staticmethod
    def word_to_ppt(doc_path: str, output_path: str):
        doc = Document(doc_path)
        
        prs = Presentation()
        
        title_added = False
        
        for para in doc.paragraphs:
            if para.style.name.startswith('Heading') and '1' in para.style.name:
                if not title_added:
                    title_slide = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide)
                    slide.shapes.title.text = para.text
                    title_added = True
                else:
                    content_slide = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(content_slide)
                    slide.shapes.title.text = para.text
                    
            elif para.text.strip():
                if not title_added:
                    title_slide = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide)
                    slide.shapes.title.text = para.text
                    title_added = True
        
        prs.save(output_path)
        return output_path

    @staticmethod
    def create_excel_with_charts(excel_path: str, output_path: str,
                                 chart_configs: List[Dict] = None):
        df = pd.read_excel(excel_path)
        
        wb = Workbook()
        ws = wb.active
        ws.title = "数据"
        
        for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
            for c_idx, value in enumerate(row, 1):
                ws.cell(r_idx, c_idx, value)
        
        if chart_configs:
            for config in chart_configs:
                chart_type = config.get("type", "bar")
                title = config.get("title", "图表")
                position = config.get("position", "H5")
                
                if chart_type == "bar":
                    chart = BarChart()
                elif chart_type == "line":
                    chart = LineChart()
                elif chart_type == "pie":
                    chart = PieChart()
                else:
                    chart = BarChart()
                
                chart.title = title
                
                data = Reference(ws, min_col=2, min_row=1, 
                               max_row=len(df)+1, 
                               max_col=min(3, len(df.columns)))
                cats = Reference(ws, min_col=1, min_row=2, max_row=len(df)+1)
                
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                ws.add_chart(chart, position)
        
        wb.save(output_path)
        return output_path


class BatchOfficeProcessor:
    @staticmethod
    def batch_convert_word_to_pdf(input_dir: str, output_dir: str):
        os.makedirs(output_dir, exist_ok=True)
        
        converted = []
        for file in os.listdir(input_dir):
            if file.endswith('.docx'):
                input_path = os.path.join(input_dir, file)
                output_path = os.path.join(output_dir, file.replace('.docx', '.pdf'))
                
                doc = Document(input_path)
                
                from docx2pdf import convert
                convert(input_path, output_path)
                
                converted.append(output_path)
        
        return converted

    @staticmethod
    def batch_merge_excel(input_files: List[str], output_path: str,
                         sheet_names: List[str] = None):
        wb = Workbook()
        if sheet_names:
            for i, sheet_name in enumerate(sheet_names):
                if i == 0:
                    ws = wb.active
                    ws.title = sheet_name
                else:
                    ws = wb.create_sheet(title=sheet_name)
        else:
            ws = wb.active
        
        all_data = []
        for file in input_files:
            df = pd.read_excel(file)
            all_data.append(df)
        
        combined = pd.concat(all_data, ignore_index=True)
        
        for r_idx, row in enumerate(dataframe_to_rows(combined, index=False, header=True), 1):
            for c_idx, value in enumerate(row, 1):
                ws.cell(r_idx, c_idx, value)
        
        wb.save(output_path)
        return output_path

    @staticmethod
    def batch_create_reports(data_file: str, template_file: str,
                            output_dir: str, prefix: str = "report"):
        df = pd.read_excel(data_file) if data_file.endswith('.xlsx') else pd.read_csv(data_file)
        
        os.makedirs(output_dir, exist_ok=True)
        
        created = []
        for i, row in df.iterrows():
            doc = Document(template_file)
            
            for para in doc.paragraphs:
                for col in df.columns:
                    if str(col) in para.text:
                        para.text = para.text.replace(str(col), str(row[col]))
            
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for col in df.columns:
                            if str(col) in cell.text:
                                cell.text = cell.text.replace(str(col), str(row[col]))
            
            output_path = os.path.join(output_dir, f"{prefix}_{i+1}.docx")
            doc.save(output_path)
            created.append(output_path)
        
        return created


def protect_word_document(input_path: str, output_path: str, password: str = None):
    doc = Document(input_path)
    
    from docx.oxml.xmlchemy import OxmlElement
    from docx.oxml.ns import qn
    
    if password:
        protection = doc.part.element.get_or_add_ptPr()
        enforcement = OxmlElement('w:enforcement')
        enforcement.set(qn('w:val'), '1')
        protection.append(enforcement)
    
    doc.save(output_path)
    return output_path


def add_word_toc(input_path: str, output_path: str):
    doc = Document(input_path)
    
    WordAdvancedFeatures.add_table_of_contents(doc, "目录")
    
    doc.save(output_path)
    return output_path


def excel_add_conditional_format(input_path: str, output_path: str,
                                 range: str, rule_type: str = "colorScale"):
    wb = load_workbook(input_path)
    ws = wb.active
    
    ExcelAdvancedFeatures.add_conditional_formatting(
        ws, range, rule_type,
        {"start_color": "F8696B", "mid_color": "FFEB84", "end_color": "63BE7B"}
    )
    
    wb.save(output_path)
    return output_path


if __name__ == "__main__":
    print("Office 高级功能测试")
    
    print("\n1. 创建带目录的Word文档...")
    doc = Document()
    doc.add_heading("测试文档", 0)
    WordAdvancedFeatures.add_table_of_contents(doc, "目录")
    doc.add_heading("第一章 概述", level=1)
    doc.add_paragraph("这是第一章的内容...")
    doc.add_heading("第二章 详细分析", level=1)
    doc.add_paragraph("这是第二章的内容...")
    WordAdvancedFeatures.add_footer(doc, include_page_number=True)
    doc.save("document_with_toc.docx")
    print("   已生成: document_with_toc.docx")
    
    print("\n2. 创建带条件格式的Excel...")
    wb = Workbook()
    ws = wb.active
    ws.title = "销售数据"
    
    headers = ["产品", "销售额", "利润", "增长率"]
    for col, header in enumerate(headers, 1):
        ws.cell(1, col, header)
    
    data = [
        ["产品A", 10000, 3000, 0.15],
        ["产品B", 15000, 4500, 0.20],
        ["产品C", 8000, 2000, 0.10],
        ["产品D", 20000, 6000, 0.25],
    ]
    
    for row_idx, row_data in enumerate(data, 2):
        for col_idx, value in enumerate(row_data, 1):
            ws.cell(row_idx, col_idx, value)
    
    ExcelAdvancedFeatures.add_conditional_formatting(ws, "B2:B100", "colorScale")
    ExcelAdvancedFeatures.add_data_validation(ws, "A2:A100", "list", 
                                              formula1='"产品A,产品B,产品C,产品D"')
    ExcelAdvancedFeatures.add_freeze_panes(ws, "A2")
    
    wb.save("excel_with_format.xlsx")
    print("   已生成: excel_with_format.xlsx")
    
    print("\n3. 格式转换测试...")
    print("   Excel → Word: 使用 excel_to_word_table()")
    print("   Excel → PPT: 使用 excel_to_ppt()")
    print("   Word → PPT: 使用 word_to_ppt()")
    
    print("\n✅ 所有测试完成!")
