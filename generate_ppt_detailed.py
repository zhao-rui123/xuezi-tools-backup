#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
零碳园区工业节能解决方案 - 详细版PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
TITLE_COLOR = RGBColor(0, 51, 102)  # 深蓝色
ACCENT_COLOR = RGBColor(0, 120, 212)  # 蓝色
SUCCESS_COLOR = RGBColor(34, 139, 34)  # 绿色
ORANGE_COLOR = RGBColor(255, 140, 0)  # 橙色
PURPLE_COLOR = RGBColor(128, 0, 128)  # 紫色
GRAY_COLOR = RGBColor(100, 100, 100)  # 灰色

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title, color=None):
    """添加章节页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg_color = color if color else TITLE_COLOR
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(4), prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_solutions_slide(prs, title, solutions, color=ACCENT_COLOR):
    """添加详细解决方案页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    y_pos = 1.2
    
    for category, items in solutions.items():
        # 类别标题（带颜色标识）
        cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(0.2), Inches(0.4))
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos + 0.1), Inches(0.2), Inches(0.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        cat_text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.4))
        tf = cat_text_box.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        y_pos += 0.5
        
        # 详细内容（两列布局）
        col1_items = items[:len(items)//2 + len(items)%2]
        col2_items = items[len(items)//2 + len(items)%2:]
        
        # 第一列
        x_pos = 0.8
        for item in col1_items:
            item_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(5.5), Inches(0.4))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✓ {item}"
            p.font.size = Pt(13)
            y_pos += 0.4
        
        # 第二列
        x_pos = 6.5
        y_pos2 = 1.7 + len(col1_items) * 0.4
        for item in col2_items:
            item_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos2), Inches(5.5), Inches(0.4))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✓ {item}"
            p.font.size = Pt(13)
            y_pos2 += 0.4
        
        y_pos = max(y_pos, y_pos2) + 0.4
    
    return slide

def add_table_slide(prs, title, headers, data):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 表格
    table = slide.shapes.add_table(len(data)+1, len(headers), Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5)).table
    
    # 设置列宽
    col_widths = [2.5, 3, 3.3, 3.5]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)
    
    # 表头
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # 数据行
    colors = [RGBColor(240, 248, 255), RGBColor(255, 255, 255)]
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.rows[row_idx+1].cells[col_idx]
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors[row_idx % 2]
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    return slide

# ========== 生成PPT ==========

# 1. 标题页
add_title_slide(prs, "工业节能解决方案", "矿山与汽车制造零碳园区")

# 2. 目录
add_section_slide(prs, "目录", ACCENT_COLOR)

# ========== 通用节能技术 ==========
add_section_slide(prs, "一、通用节能技术", SUCCESS_COLOR)

add_solutions_slide(prs, "1. 能源管理系统(EMS)", {
    "系统功能": [
        "实时能耗监测与数据分析",
        "多能互补智能调度",
        "碳排放核算与管理",
        "能耗异常预警",
        "能效对标分析",
        "移动端远程监控"
    ]
}, SUCCESS_COLOR)

add_solutions_slide(prs, "2. 设备变频改造(VFD)", {
    "改造范围": [
        "电机变频驱动",
        "水泵变频调速",
        "空压机变频控制",
        "风机变频调节",
        "空调系统变频",
        "传送带变频驱动"
    ],
    "节能效果": [
        "综合节能15%-30%",
        "降低设备启动电流",
        "延长设备寿命",
        "减少维护成本",
        "提高功率因数"
    ]
}, SUCCESS_COLOR)

add_solutions_slide(prs, "3. 余热回收利用", {
    "回收方式": [
        "烟气余热回收",
        "废水余热回收",
        "设备散热回收",
        "工艺余热回收",
        "余热锅炉利用",
        "余热制冷空调"
    ],
    "应用场景": [
        "预热助燃空气",
        "生产生活热水",
        "冬季供暖",
        "吸收式制冷",
        "余热发电"
    ]
}, SUCCESS_COLOR)

add_solutions_slide(prs, "4. 智能照明系统", {
    "技术措施": [
        "LED灯具替换",
        "光照度传感器",
        "人体感应控制",
        "时间场景控制",
        "光照分区调节",
        "智能照明管理系统"
    ],
    "节能效果": [
        "节能50%-70%",
        "自动调光补偿",
        "故障自动报警",
        "寿命长达5万小时"
    ]
}, SUCCESS_COLOR)

add_solutions_slide(prs, "5. 建筑节能", {
    "围护结构": [
        "外墙保温系统",
        "节能门窗",
        "屋面隔热防水",
        "外墙反射隔热"
    ],
    "空调系统": [
        "高效节能机组",
        "智能温控系统",
        "新风热回收",
        "分区独立控制"
    ]
}, SUCCESS_COLOR)

# ========== 汽车制造专项 ==========
add_section_slide(prs, "二、汽车制造节能方案", ACCENT_COLOR)

add_solutions_slide(prs, "1. 工业炉窑节能", {
    "蓄热式燃烧技术": [
        "高温空气燃烧(HTAC)",
        "蓄热式烧嘴(Regen)",
        "脉冲燃烧控制",
        "炉温均匀性优化"
    ],
    "余热回收": [
        "烟气余热预热空气",
        "余热锅炉产蒸汽",
        "余热回收用于涂装前处理"
    ],
    "智能控制": [
        "分区温度精确控制",
        "炉压自动调节",
        "燃烧优化算法",
        "故障诊断预警"
    ]
}, ACCENT_COLOR)

add_solutions_slide(prs, "2. 涂装车间节能", {
    "废气余热回收": [
        "RTO废热回收",
        "涂装废气余热利用",
        "烘干室余热回收"
    ],
    "工艺优化": [
        "水性漆替代溶剂漆",
        "紧凑型涂装工艺",
        "机器人自动喷涂",
        "高效换色系统"
    ],
    "空调节能": [
        "温湿度独立控制",
        "转轮除湿系统",
        "空调变频控制",
        "局部送风技术"
    ]
}, ACCENT_COLOR)

add_solutions_slide(prs, "3. 空压站节能", {
    "设备升级": [
        "永磁变频空压机",
        "高效螺杆空压机",
        "离心式空压机",
        "无油润滑空压机"
    ],
    "系统优化": [
        "空压机群联控",
        "泄漏检测与修复",
        "压降优化",
        "余热回收利用"
    ],
    "智能控制": [
        "负荷预测控制",
        "恒压供气",
        "远程监控系统",
        "能效统计分析"
    ]
}, ACCENT_COLOR)

add_solutions_slide(prs, "4. 铸造车间节能", {
    "熔炼设备": [
        "感应电炉取代冲天炉",
        "保温熔炼炉",
        "高效浇注系统",
        "炉料预热"
    ],
    "砂处理": [
        "旧砂冷却再生",
        "混砂机变频改造",
        "除尘系统优化"
    ],
    "工艺优化": [
        "精密铸造工艺",
        "消失模铸造",
        "低压铸造",
        "离心铸造"
    ]
}, ACCENT_COLOR)

add_solutions_slide(prs, "5. 机加工节能", {
    "设备升级": [
        "高效数控机床",
        "电主轴技术",
        "直线电机驱动",
        "高速加工中心"
    ],
    "工艺优化": [
        "干式切削技术",
        "微量润滑(MQL)",
        "高效刀具",
        "加工参数优化"
    ],
    "辅助系统": [
        "切削液净化回收",
        "机床照明LED",
        "排屑输送节能"
    ]
}, ACCENT_COLOR)

add_solutions_slide(prs, "6. 物流与仓储", {
    "搬运设备": [
        "电动叉车替换",
        "AGV无人搬运",
        "输送线变频",
        "提升机能量回收"
    ],
    "仓储系统": [
        "自动化立体仓库",
        "智能分拣系统",
        "仓储节能照明",
        "温湿度监控"
    ]
}, ACCENT_COLOR)

# ========== 矿山专项 ==========
add_section_slide(prs, "三、矿山节能方案", ORANGE_COLOR)

add_solutions_slide(prs, "1. 提升系统节能", {
    "设备升级": [
        "直驱式提升机",
        "变频调速驱动",
        "永磁同步电机"
    ],
    "能量回收": [
        "势能回收发电",
        "制动能量回馈",
        "超级电容储能"
    ],
    "智能控制": [
        "最优升降曲线",
        "智能称重装载",
        "无人值守控制"
    ]
}, ORANGE_COLOR)

add_solutions_slide(prs, "2. 空压站节能", {
    "设备升级": [
        "离心变频空压机",
        "永磁变频螺杆机",
        "高效油气分离"
    ],
    "系统优化": [
        "余热回收利用",
        "泄漏智能检测",
        "压风管网优化",
        "集中远程控制"
    ],
    "按需供气": [
        "分段压力控制",
        "末端储气罐",
        "用气需求预测"
    ]
}, ORANGE_COLOR)

add_solutions_slide(prs, "3. 通风系统节能", {
    "智能通风": [
        "按需供风控制",
        "风量自动调节",
        "风流组织优化",
        "瓦斯智能排放"
    ],
    "设备升级": [
        "高效对旋风机",
        "变频调速驱动",
        "消声隔音措施"
    ],
    "自然通风": [
        "竖井自然通风",
        "巷道风流优化",
        "通风阻力降低"
    ]
}, ORANGE_COLOR)

add_solutions_slide(prs, "4. 排水系统节能", {
    "设备升级": [
        "高效潜水泵",
        "变频调速控制",
        "永磁同步电机"
    ],
    "智能控制": [
        "水位智能监控",
        "排水量预测",
        "峰谷电价运行",
        "无人值守控制"
    ],
    "系统优化": [
        "管网泄漏检测",
        "阀门优化调节",
        "水仓有效利用"
    ]
}, ORANGE_COLOR)

add_solutions_slide(prs, "5. 采矿设备电动化", {
    "运输设备": [
        "电动宽体车",
        "电机车直流改交流",
        "胶带输送机变频"
    ],
    "采掘设备": [
        "电动挖掘机",
        "液压站电驱动",
        "凿岩台车电动化"
    ],
    "辅助设备": [
        "电动无轨胶轮车",
        "架空乘人装置",
        "照明系统LED"
    ]
}, ORANGE_COLOR)

add_solutions_slide(prs, "6. 选矿节能", {
    "破碎筛分": [
        "高压辊磨机",
        "高效圆锥破碎机",
        "智能筛分控制"
    ],
    "磨矿分级": [
        "球磨机优化衬板",
        "磨矿浓度优化",
        "分级效率提升"
    ],
    "浮选浓缩": [
        "浮选机高效叶轮",
        "药剂用量优化",
        "浓缩池高效刮泥"
    ]
}, ORANGE_COLOR)

# ========== 可再生能源 ==========
add_section_slide(prs, "四、可再生能源利用", PURPLE_COLOR)

add_solutions_slide(prs, "1. 光伏发电", {
    "分布式光伏": [
        "厂房屋顶光伏",
        "停车场光伏车棚",
        "空地集中式",
        "幕墙一体化"
    ],
    "储能配套": [
        "光伏储能系统",
        "削峰填谷",
        "应急备电"
    ]
}, PURPLE_COLOR)

add_solutions_slide(prs, "2. 储能系统", {
    "电力储能": [
        "电池储能系统(BESS)",
        "飞轮储能",
        "压缩空气储能"
    ],
    "应用场景": [
        "需量管理",
        "峰谷套利",
        "新能源消纳",
        "电网调频"
    ]
}, PURPLE_COLOR)

add_solutions_slide(prs, "3. 绿色电力", {
    "绿电采购": [
        "风电采购",
        "光伏采购",
        "绿电证书"
    ],
    "碳减排": [
        "碳排放核算",
        "碳交易",
        "碳中和规划"
    ]
}, PURPLE_COLOR)

# ========== 投资与效益 ==========
add_table_slide(prs, "节能投资与效益估算", 
    ["节能措施", "投资估算", "年节能效益", "回收期"],
    [
        ["能源管理系统(EMS)", "50-100万元", "10-20万元", "3-5年"],
        ["变频改造", "30-80万元", "15-30万元", "1.5-3年"],
        ["余热回收", "100-300万元", "30-80万元", "3-5年"],
        ["LED照明", "20-50万元", "10-20万元", "1.5-2年"],
        ["空压站节能", "50-150万元", "20-50万元", "2-3年"],
        ["光伏发电", "200-500万元", "40-100万元", "4-6年"],
        ["储能系统", "300-800万元", "30-60万元", "8-12年"]
    ]
)

# ========== 实施路径 ==========
add_solutions_slide(prs, "节能实施路径", {
    "第一阶段(1-6月)": [
        "能源审计与诊断",
        "建立能效基准",
        "Quick-win项目实施",
        "EMS系统部署"
    ],
    "第二阶段(6-18月)": [
        "设备升级改造",
        "工艺优化调整",
        "系统智能控制",
        "人员培训推广"
    ],
    "第三阶段(18-36月)": [
        "可再生能源配套",
        "碳管理体系建设",
        "持续改进优化",
        "行业标杆打造"
    ]
}, SUCCESS_COLOR)

# 总结页
add_section_slide(prs, "谢谢", TITLE_COLOR)

# 保存文件
output_path = os.path.expanduser("~/.openclaw/workspace/零碳园区节能解决方案_详细版.pptx")
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
