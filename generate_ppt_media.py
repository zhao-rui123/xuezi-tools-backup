#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
公用介质用能设备明细PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0, 51, 102)
ACCENT_COLOR = RGBColor(0, 120, 212)
ORANGE_COLOR = RGBColor(255, 140, 0)
SUCCESS_COLOR = RGBColor(34, 139, 34)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_media_detail_slide(prs, title, media_data, color=ACCENT_COLOR):
    """添加公用介质详细设备页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # 绘制表格
    y_pos = 1.1
    
    for media, info in media_data.items():
        # 介质名称和占比
        media_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        tf = media_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{media}\n{info['ratio']}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        # 设备列表
        devices = info['devices']
        col1 = devices[:len(devices)//2 + len(devices)%2]
        col2 = devices[len(devices)//2 + len(devices)%2:]
        
        # 第一列
        x_pos = 2.8
        for dev in col1:
            dev_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(4.5), Inches(0.4))
            tf = dev_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {dev}"
            p.font.size = Pt(13)
            y_pos += 0.4
        
        # 第二列
        x_pos = 7.5
        y_pos2 = 1.1 + len(col1) * 0.4
        for dev in col2:
            dev_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos2), Inches(4.5), Inches(0.4))
            tf = dev_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {dev}"
            p.font.size = Pt(13)
            y_pos2 += 0.4
        
        y_pos = max(y_pos, y_pos2) + 0.5
    
    return slide

# 1. 标题页
add_title_slide(prs, "公用介质用能设备明细", "矿山与汽车制造零碳园区")

# 2. 汽车制造 - 公用介质设备
add_media_detail_slide(prs, "汽车制造零碳园区 - 公用介质用能设备明细", {
    "电力": {
        "ratio": "80%",
        "devices": [
            "工业炉窑（热处理、铸造）",
            "空压机站",
            "各类电机设备",
            "空调与暖通系统",
            "数控机床",
            "工业机器人",
            "涂装生产线",
            "输送设备",
            "照明系统",
            "检测设备"
        ]
    },
    "天然气": {
        "ratio": "15%",
        "devices": [
            "工业炉窑（加热炉）",
            "锅炉（蒸汽/热水）",
            "热处理炉",
            "涂装烘干炉",
            "钎焊设备",
            "厨房餐饮设备"
        ]
    },
    "蒸汽": {
        "ratio": "5%",
        "devices": [
            "工艺加热",
            "前处理清洗",
            "涂装脱脂",
            "空调供暖",
            "设备消毒",
            "实验室用气"
        ]
    }
}, ACCENT_COLOR)

# 3. 矿山 - 公用介质设备
add_media_detail_slide(prs, "矿山零碳园区 - 公用介质用能设备明细", {
    "电力": {
        "ratio": "91%",
        "devices": [
            "提升机",
            "空压机",
            "排水泵",
            "通风机",
            "各类电机",
            "照明系统",
            "输送胶带机",
            "破碎筛分设备",
            "浮选设备",
            "压滤设备",
            "除尘设备",
            "监控通信设备"
        ]
    },
    "柴油": {
        "ratio": "7%",
        "devices": [
            "电动宽体车",
            "挖掘机",
            "装载机",
            "推土机",
            "钻探设备",
            "备用发电机",
            "应急照明"
        ]
    },
    "天然气": {
        "ratio": "2%",
        "devices": [
            "锅炉房",
            "冬季取暖",
            "员工厨房",
            "办公楼供暖"
        ]
    }
}, ORANGE_COLOR)

# 4. 汇总对比
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "公用介质用能对比汇总"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR

# 汇总表格
table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5)).table

headers = ["公用介质", "汽车制造占比", "矿山占比", "主要差异"]
for i, h in enumerate(headers):
    cell = table.rows[0].cells[i]
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = TITLE_COLOR
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    cell.text_frame.paragraphs[0].font.bold = True

data = [
    ["电力", "80%", "91%", "矿山更高，提升/排水/通风为主"],
    ["天然气", "15%", "2%", "汽车制造用于热加工"],
    ["蒸汽", "5%", "0%", "汽车制造有工艺蒸汽需求"],
    ["柴油", "0%", "7%", "矿山运输设备为主"]
]

for row_idx, row_data in enumerate(data):
    for col_idx, value in enumerate(row_data):
        cell = table.rows[row_idx+1].cells[col_idx]
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
        cell.text_frame.paragraphs[0].font.size = Pt(12)

# 保存
output_path = os.path.expanduser("~/.openclaw/workspace/公用介质用能设备明细.pptx")
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
