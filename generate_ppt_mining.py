#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
矿山零碳园区解决方案PPT - 完整版
按照模板格式：项目概述、用能特征分析、技术方案设计、碳排放特征、零碳路径规划、效益分析
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0, 51, 102)
ACCENT_COLOR = RGBColor(0, 120, 212)
ORANGE_COLOR = RGBColor(255, 140, 0)
GREEN_COLOR = RGBColor(34, 139, 34)

def add_cover_slide(prs, title, subtitle, industry):
    """封面页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部色带
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE_COLOR
    shape.line.fill.background()
    
    # 标题背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), prs.slide_width, Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 行业标签
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.6))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📋 {industry}"
    p.font.size = Pt(20)
    p.font.color.rgb = ORANGE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, section_title, color=None):
    """章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg_color = color if color else TITLE_COLOR
    
    # 左侧色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    
    # 章节号
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(3.5), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(5.2), Inches(3), Inches(7.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = bg_color
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_full_content_slide(prs, title, sections):
    """全内容页 - 多个section，每个有标题和要点"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    y_pos = 1.0
    for section_title, items in sections.items():
        # 小节标题
        title_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.4))
        tf = title_box2.text_frame
        p = tf.paragraphs[0]
        p.text = f"◆ {section_title}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ORANGE_COLOR
        y_pos += 0.4
        
        # 要点
        for item in items:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y_pos + 0.06), Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT_COLOR
            dot.line.fill.background()
            
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.32))
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(11)
            y_pos += 0.32
        
        y_pos += 0.15
    
    return slide

def add_table_slide(prs, title, headers, data):
    """表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 表格
    num_rows = len(data) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.3), Inches(1.0), Inches(12.7), Inches(0.5)).table
    
    # 表头
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # 数据行
    colors = [RGBColor(255, 245, 235), RGBColor(255, 255, 255)]
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.rows[row_idx+1].cells[col_idx]
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors[row_idx % 2]
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    return slide

def add_timeline_slide(prs, title, phases):
    """时间轴页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 时间轴线
    y_line = Inches(3.8)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y_line, Inches(11.7), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE_COLOR
    line.line.fill.background()
    
    # 各阶段
    phase_width = 12.0 / len(phases)
    colors_phase = [ACCENT_COLOR, ORANGE_COLOR, GREEN_COLOR]
    
    for i, (phase_name, goals) in enumerate(phases):
        x_pos = 0.8 + i * phase_width
        
        # 节点圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + phase_width/2 - 0.15), Inches(3.65), Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors_phase[i % 3]
        dot.line.fill.background()
        
        # 阶段名称
        name_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.3), Inches(phase_width), Inches(0.6))
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase_name
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = colors_phase[i % 3]
        p.alignment = PP_ALIGN.CENTER
        
        # 目标内容
        y_goal = 2.0
        for goal in goals:
            goal_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(y_goal), Inches(phase_width - 0.2), Inches(0.4))
            tf = goal_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {goal}"
            p.font.size = Pt(11)
            y_goal += 0.4
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5.8), Inches(0.4))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"◆ {left_title}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE_COLOR
    
    # 左栏内容
    y_pos = 1.5
    for item in left_items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos + 0.06), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_COLOR
        dot.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(5.5), Inches(0.4))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        y_pos += 0.4
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.0), Inches(5.8), Inches(0.4))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"◆ {right_title}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN_COLOR
    
    # 右栏内容
    y_pos = 1.5
    for item in right_items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(y_pos + 0.06), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN_COLOR
        dot.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(7.2), Inches(y_pos), Inches(5.5), Inches(0.4))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        y_pos += 0.4
    
    return slide

def add_benefit_slide(prs, title, benefits):
    """效益分析页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.15), prs.slide_width, Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 三列效益
    col_width = 4.0
    colors = [ACCENT_COLOR, ORANGE_COLOR, GREEN_COLOR]
    icons = ["💰", "🏭", "🌱"]
    
    for i, (benefit_type, items) in enumerate(benefits.items()):
        x_pos = 0.5 + i * col_width
        
        # 顶部图标和标题
        icon_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.0), Inches(col_width), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icons[i]} {benefit_type}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER
        
        # 分隔线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_pos + 0.5), Inches(1.6), Inches(col_width - 1), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = colors[i]
        line.line.fill.background()
        
        # 内容
        y_pos = 1.8
        for item in items:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 0.2), Inches(y_pos + 0.06), Inches(0.08), Inches(0.08))
            dot.fill.solid()
            dot.fill.fore_color.rgb = colors[i]
            dot.line.fill.background()
            
            text_box = slide.shapes.add_textbox(Inches(x_pos + 0.4), Inches(y_pos), Inches(col_width - 0.5), Inches(0.45))
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(12)
            y_pos += 0.5
    
    return slide

# ========== 封面 ==========
add_cover_slide(prs, "矿山零碳园区解决方案", "Green Mining Park Solution", "矿山行业零碳园区")

# ========== 01 项目概述 ==========
add_section_slide(prs, "01", "项目概述", ORANGE_COLOR)

add_full_content_slide(prs, "项目建设背景", {
    "行业背景": [
        "矿业是我国国民经济的基础性产业，能源消耗大，碳排放强度高",
        "矿山作业涉及采矿、选矿、提升、运输、排水、通风等高能耗工序",
        "矿山企业面临碳排放约束和绿色转型的双重压力"
    ],
    "数据特征": [
        "公用介质用能占比：电力91%、柴油7%、天然气2%",
        "主要工艺用能：采矿41%、选矿38%、提升运输21%",
        "主要设备用能：提升机30%、空压机25%、水泵18%、通风机15%、电机12%"
    ],
    "政策要求": [
        "双碳目标：2030年碳达峰、2060年碳中和",
        "能耗双控：严格控制能耗强度和总量",
        "绿色矿山：2025年基本实现全部省级绿色矿山"
    ]
})

add_full_content_slide(prs, "项目建设目标", {
    "近期目标(2023-2025)": [
        "能效提升25%，单位产品能耗下降",
        "碳排放强度下降18%",
        "建成省级绿色矿山"
    ],
    "中期目标(2026-2030)": [
        "绿电占比达到40%",
        "碳排放强度下降45%",
        "建成国家级绿色矿山示范"
    ],
    "远期目标(2031-2035)": [
        "实现碳中和",
        "100%使用绿色电力",
        "建成零碳智慧矿山标杆"
    ]
})

# ========== 02 用能特征分析 ==========
add_section_slide(prs, "02", "用能特征分析", ORANGE_COLOR)

add_table_slide(prs, "主要工艺用能占比", 
    ["工艺环节", "能耗占比", "主要耗能设备", "节能重点"],
    [
        ["采矿", "41%", "凿岩机、铲运机、矿车", "设备电动化、智能化"],
        ["选矿", "38%", "破碎机、球磨机、浮选机", "高效设备、余热回收"],
        ["提升运输", "21%", "提升机、胶带机、电机车", "能量回收、变频调速"]
    ]
)

add_table_slide(prs, "主要设备用能占比",
    ["主要设备", "用能占比", "设备数量", "节能改造方向"],
    [
        ["提升机", "30%", "主井/副井提升机", "变频调速+能量回收"],
        ["空压机", "25%", "空压机站", "变频+余热利用+群控"],
        ["水泵", "18%", "井下排水泵", "变频调速+峰谷运行"],
        ["通风机", "15%", "主扇/局扇", "智能通风+变频"],
        ["其他电机", "12%", "各类辅助电机", "变频改造+能效提升"]
    ]
)

add_table_slide(prs, "公用介质用能占比",
    ["介质类型", "用能占比", "主要用途", "替代方向"],
    [
        ["电力", "91%", "提升、排水、通风、照明", "绿电替代+节能降耗"],
        ["柴油", "7%", "运输车辆、工程机械", "电动化替代"],
        ["天然气", "2%", "锅炉取暖、员工设施", "热泵替代"]
    ]
)

add_timeline_slide(prs, "用能结构优化趋势",
    [
        ("现状\n2023", ["电力91%", "柴油7%", "天然气2%", "碳强度1.2"]),
        ("近期\n2025", ["电力85%", "柴油5%", "天然气1%", "绿电15%"]),
        ("中期\n2030", ["电力75%", "柴油3%", "天然气0%", "绿电40%"]),
        ("远期\n2035", ["电力60%", "柴油1%", "天然气0%", "绿电100%"])
    ]
)

# ========== 03 技术方案设计 ==========
add_section_slide(prs, "03", "技术方案设计", ORANGE_COLOR)

add_two_column_slide(prs, "技术方案设计 - 能源替代与设备节能",
    "能源替代方案", [
        "分布式光伏：办公区/生活区屋顶建设",
        "储能系统：配置储能平抑波动、削峰填谷",
        "电动化替代：电动宽体车替代柴油车辆",
        "绿电采购：购买风电/光伏绿证"
    ],
    "设备节能方案", [
        "提升机：变频调速+势能回收发电",
        "空压站：变频+余热利用+按需供气",
        "通风系统：智能调控+变频调速",
        "排水系统：变频调速+峰谷电价运行"
    ]
)

add_two_column_slide(prs, "技术方案设计 - 智能化管理与能源结构优化",
    "智能化管理", [
        "矿山智慧能源平台建设",
        "设备状态在线监测诊断",
        "能耗实时监测与分析",
        "碳排放统计与核算"
    ],
    "能源结构优化", [
        "电力占比：91%→60%",
        "光伏装机：新增X MWp",
        "储能配置：X MWh",
        "绿电采购：年购绿证X万份"
    ]
)

add_full_content_slide(prs, "重点节能技术方案", {
    "提升系统节能": [
        "采用变频调速技术，根据负载自动调节提升速度",
        "配置能量回收系统，将势能转化为电能回馈电网",
        "采用永磁同步电机，提高电机效率",
        "智能优化升降曲线，减少无效能耗"
    ],
    "空压站节能": [
        "采用永磁变频空压机，替代传统工频空压机",
        "配置余热回收装置，利用压缩热制备热水/供暖",
        "实施空压机群联控，按需供气，减少卸载能耗",
        "智能泄漏检测，及时发现并修复压缩空气泄漏"
    ],
    "通风系统节能": [
        "实施智能通风控制，根据瓦斯浓度和作业区域按需供风",
        "采用高效对旋风机，替代老旧风机",
        "变频调速控制，风量实时调节",
        "优化巷道通风网络，降低通风阻力"
    ]
})

# ========== 04 碳排放特征 ==========
add_section_slide(prs, "04", "碳排放特征", ORANGE_COLOR)

add_table_slide(prs, "碳排放结构分析",
    ["排放源", "碳排放占比", "主要来源", "减排技术路径"],
    [
        ["电力消耗", "75%", "提升/排水/通风/照明", "绿电替代+设备节能"],
        ["柴油消耗", "20%", "运输车辆/工程机械", "电动化替代"],
        ["天然气燃烧", "3%", "锅炉/取暖", "热泵替代"],
        ["其他排放", "2%", "炸药/耗材等", "工艺优化改进"]
    ]
)

add_full_content_slide(prs, "碳排放强度与减排潜力", {
    "碳排放强度现状": [
        "当前碳排放强度：约 1.2 tCO2/万吨矿产品",
        "年碳排放总量：约 X 万吨CO2",
        "万元产值碳排放：约 0.9 tCO2/万元"
    ],
    "减排潜力分析": [
        "设备节能降耗：可减排 30%（变频改造+能量回收）",
        "绿电替代：可减排 25%（光伏+储能+绿证采购）",
        "电动化替代：可减排 20%（电动车辆替代柴油）",
        "智能化管理：可减排 10%（按需供能+精细管理）"
    ],
    "减排目标": [
        "2025年：碳排放强度下降18%，降至约0.98 tCO2/吨",
        "2030年：碳排放强度下降45%，降至约0.66 tCO2/吨",
        "2035年：实现碳中和，碳排放强度趋近于零"
    ]
})

add_full_content_slide(prs, "碳核算与碳管理", {
    "碳核算体系": [
        "建立矿山碳排放核算方法学",
        "涵盖范围一（直接排放）、范围二（电力间接排放）",
        "建立碳排放因子数据库",
        "实现碳排放数据实时监测与统计"
    ],
    "碳管理体系": [
        "建立企业碳排放管理制度",
        "配置专职碳排放管理人员",
        "制定年度碳排放预算",
        "开展碳排放绩效评价"
    ],
    "碳交易准备": [
        "研究碳市场交易规则",
        "评估碳资产开发潜力",
        "建立碳资产管理台账",
        "探索CCER项目开发"
    ]
})

# ========== 05 零碳路径规划 ==========
add_section_slide(prs, "05", "零碳路径规划", ORANGE_COLOR)

add_timeline_slide(prs, "零碳路径规划 - 三阶段实施路线",
    [
        ("近期\n2023-2025", ["能效提升25%", "设备变频改造", "EMS系统上线", "省级绿色矿山", "碳强度-18%"]),
        ("中期\n2026-2030", ["绿电占比40%", "光伏+储能", "电动宽体车", "智能矿山", "碳强度-45%"]),
        ("远期\n2031-2035", ["碳中和", "100%绿电", "全面电动化", "零碳智慧", "碳汇抵消"])
    ]
)

add_full_content_slide(prs, "阶段性重点任务", {
    "近期重点(2023-2025)": [
        "完成能效诊断，制定节能改造方案",
        "实施提升机、排水泵、通风机变频改造",
        "建设能源管理系统(EMS)平台",
        "建成省级绿色矿山"
    ],
    "中期重点(2026-2030)": [
        "建设分布式光伏+储能系统",
        "推进运输车辆电动化替代",
        "建设智慧能源管控平台",
        "绿电采购比例达到40%"
    ],
    "远期重点(2031-2035)": [
        "实现100%绿电消费",
        "全面电动化替代柴油",
        "建设零碳智慧矿山",
        "通过碳汇实现碳中和"
    ]
})

add_table_slide(prs, "投资估算与资金来源",
    ["建设内容", "投资估算", "资金来源", "实施周期"],
    [
        ["设备变频改造", "X万元", "企业自筹+节能贷款", "1-2年"],
        ["能源管理系统", "X万元", "企业自筹", "1年"],
        ["分布式光伏", "X万元", "企业自筹+绿金贷款", "2-3年"],
        ["储能系统", "X万元", "企业自筹+政策补贴", "1-2年"],
        ["电动化替代", "X万元", "企业自筹+设备租赁", "3-5年"],
        ["智慧平台", "X万元", "企业自筹", "2年"]
    ]
)

# ========== 06 效益分析 ==========
add_section_slide(prs, "06", "效益分析", ORANGE_COLOR)

add_benefit_slide(prs, "效益分析 - 经济/社会/环境效益",
    {
        "经济效益": [
            "年节能收益约X万元",
            "谷电利用降低成本",
            "设备维护成本降低",
            "碳交易潜在收益",
            "政策补贴收入"
        ],
        "社会效益": [
            "改善井下作业环境",
            "提升安全生产水平",
            "提供绿色就业岗位",
            "行业示范引领作用",
            "助力地方经济发展"
        ],
        "环境效益": [
            "年减碳X万吨CO2",
            "粉尘排放大幅减少",
            "噪音污染有效降低",
            "生态恢复治理",
            "助力双碳目标实现"
        ]
    }
)

add_full_content_slide(prs, "综合效益测算", {
    "节能效益": [
        "年节约标煤：X吨",
        "年节电：X万kWh",
        "年节省运行成本：X万元"
    ],
    "减排效益": [
        "年减少CO2排放：X万吨",
        "年减少柴油消耗：X吨",
        "年减少粉尘排放：X吨"
    ],
    "经济评价": [
        "项目总投资：X万元",
        "年运行收益：X万元",
        "投资回收期：X年",
        "内部收益率：X%"
    ]
})

# ========== 总结页 ==========
add_section_slide(prs, "谢谢", "THANK YOU", TITLE_COLOR)

# 保存文件
output_path = os.path.expanduser("~/.openclaw/workspace/矿山零碳园区解决方案.pptx")
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
